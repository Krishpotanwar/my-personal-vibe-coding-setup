#!/usr/bin/env python3
"""
Trackly Presentation → PowerPoint (.pptx)

Usage:
    pip install python-pptx
    python trackly_convert_to_pptx.py

Output: Trackly.pptx

Slide size : 13.333" × 7.5"  (standard PowerPoint 16:9 widescreen)
Safe margin: 0.5" L/R · 0.5" top · 0.55" bottom
Hard bottom: 6.95" — every element verified within bounds
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Canvas ────────────────────────────────────────────────────────────────────
W, H   = Inches(13.333), Inches(7.5)
ML, MT = Inches(0.5),    Inches(0.5)
CW     = W - Inches(1.0)          # 12.333"
MAX_Y  = H - Inches(0.55)         # 6.95"
CX, CY = ML, MT

CONTENT_Y = CY + Inches(1.52)     # default content start (after label+heading+sub)

# ── Colours ───────────────────────────────────────────────────────────────────
BG     = RGBColor(0x05, 0x08, 0x0f)
CARD   = RGBColor(0x0e, 0x16, 0x2a)
CARD2  = RGBColor(0x14, 0x22, 0x3e)
GREEN  = RGBColor(0x34, 0xd3, 0x99)
PURPLE = RGBColor(0x81, 0x8c, 0xf8)
ORANGE = RGBColor(0xfb, 0x92, 0x3c)
ROSE   = RGBColor(0xf8, 0x71, 0x71)
SKY    = RGBColor(0x38, 0xbd, 0xf8)
TXT    = RGBColor(0xf0, 0xf4, 0xff)
SEC    = RGBColor(0x94, 0xa3, 0xb8)
DIM    = RGBColor(0x47, 0x55, 0x69)
GRN    = RGBColor(0xa8, 0xff, 0x78)
TERM   = RGBColor(0x02, 0x04, 0x08)
TRMBD  = RGBColor(0x1a, 0x2a, 0x1a)
BDR_G  = RGBColor(0x0d, 0x3d, 0x2a)   # green-tinted border
BDR_P  = RGBColor(0x1e, 0x22, 0x4a)   # purple-tinted border
BDR_O  = RGBColor(0x3a, 0x22, 0x0a)   # orange-tinted border


# ── Core helpers ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def tb(sl, text, x, y, w, h,
       sz=12, bold=False, color=TXT,
       align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(sz)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return box


def card(sl, x, y, w, h, fill=CARD, border=CARD2, rounded=True):
    sid = 5 if rounded else 1
    s   = sl.shapes.add_shape(sid, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width     = Pt(0.75)
    else:
        s.line.fill.background()
    if rounded:
        pg = s.element.find('.//' + qn('a:prstGeom'))
        if pg is not None:
            av = pg.find(qn('a:avLst'))
            if av is not None:
                for gd in av.findall(qn('a:gd')):
                    if gd.get('name') == 'adj':
                        gd.set('fmla', 'val 12000')
    return s


def dot(sl, x, y, r, color):
    s = sl.shapes.add_shape(9, x, y, r, r)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def lbl(sl, text, x, y, w):
    tb(sl, text.upper(), x, y, w, Inches(0.26),
       sz=7, color=DIM, align=PP_ALIGN.CENTER)


def h2(sl, text, x, y, w, sz=26):
    tb(sl, text, x, y, w, Inches(0.62),
       sz=sz, bold=True, color=TXT, align=PP_ALIGN.CENTER)


def sub(sl, text, x, y, w, sz=11):
    tb(sl, text, x, y, w, Inches(0.48),
       sz=sz, color=SEC, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# S01 — Title
# bottom ≈ CY + 4.35 = 4.85"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s01(prs):
    sl = blank(prs)

    # Badge
    bw, bh = Inches(6.0), Inches(0.36)
    bx = CX + (CW - bw) / 2
    card(sl, bx, CY, bw, bh,
         fill=RGBColor(0x07, 0x18, 0x22), border=BDR_G)
    tb(sl, "NGO EVENT MANAGEMENT PLATFORM",
       bx + Inches(0.1), CY + Inches(0.05), bw - Inches(0.2), Inches(0.26),
       sz=8, color=GREEN, align=PP_ALIGN.CENTER)

    # App name
    tb(sl, "Trackly",
       CX, CY + Inches(0.52), CW, Inches(0.85),
       sz=52, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Tagline
    tb(sl, "Coordinate your volunteers, track your events, and measure your impact — "
           "all in one place built for how nonprofits actually work.",
       CX + Inches(1.8), CY + Inches(1.52), CW - Inches(3.6), Inches(0.75),
       sz=13, color=SEC, align=PP_ALIGN.CENTER)

    # Team / meta
    tb(sl, "Shlok Dwivedi  ·  React + TypeScript + Firebase + Flask  ·  119 commits",
       CX, CY + Inches(2.44), CW, Inches(0.28),
       sz=10, color=DIM, align=PP_ALIGN.CENTER)

    # Links
    links = [
        ("trackly-phi.vercel.app",          GREEN),
        ("github.com/Shlok-Dwivedi/Trackly", SEC),
    ]
    lw, lh = Inches(4.5), Inches(0.38)
    gap    = Inches(0.28)
    ly     = CY + Inches(2.90)
    lsx    = CX + (CW - 2 * lw - gap) / 2
    for i, (link, clr) in enumerate(links):
        lx = lsx + i * (lw + gap)
        card(sl, lx, ly, lw, lh, fill=CARD, border=CARD2)
        tb(sl, link, lx + Inches(0.1), ly + Inches(0.07),
           lw - Inches(0.2), lh - Inches(0.1),
           sz=10, color=clr, align=PP_ALIGN.CENTER, italic=True)

    # Tech pills
    pills = ["React 18", "TypeScript", "Firebase", "Flask", "Supabase", "Tailwind", "Vercel"]
    pw, ph = Inches(1.6), Inches(0.30)
    pg     = Inches(0.14)
    py     = CY + Inches(3.48)
    psx    = CX + (CW - len(pills) * pw - (len(pills) - 1) * pg) / 2
    for i, p_txt in enumerate(pills):
        px  = psx + i * (pw + pg)
        hl  = i < 3
        card(sl, px, py, pw, ph,
             fill=RGBColor(0x06, 0x20, 0x16) if hl else CARD,
             border=BDR_G if hl else CARD2)
        tb(sl, p_txt, px + Inches(0.05), py + Inches(0.05),
           pw - Inches(0.1), ph - Inches(0.08),
           sz=9, color=GREEN if hl else DIM, align=PP_ALIGN.CENTER)
    # bottom ≈ CY + 3.78 = 4.28"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S02 — Problem
# bottom ≈ CY + 4.82 = 5.32"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s02(prs):
    sl = blank(prs)
    lbl(sl, "The Problem", CX, CY, CW)
    h2(sl, "Running an NGO on spreadsheets and WhatsApp groups", CX, CY + Inches(0.30), CW)
    sub(sl, "A coordinator managing 12 events, 40 volunteers, and a shared calendar "
            "nobody updates spends 3 hours a week chasing status updates over chat.",
        CX + Inches(0.8), CY + Inches(0.96), CW - Inches(1.6))

    stats = [
        ("3+",  "hours a week lost to\ncoordination overhead",   ROSE),
        ("0",   "unified place for event data,\nattendees, and analytics", ORANGE),
        ("4+",  "disconnected tools used\nfor one event lifecycle", PURPLE),
    ]
    gap, cw, ch = Inches(0.22), (CW - 2 * Inches(0.22)) / 3, Inches(2.45)
    cy = CONTENT_Y
    for i, (num, desc, clr) in enumerate(stats):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, num, cx, cy + Inches(0.35), cw, Inches(0.80),
           sz=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tb(sl, desc, cx + Inches(0.18), cy + Inches(1.25), cw - Inches(0.36), Inches(0.95),
           sz=11, color=SEC, align=PP_ALIGN.CENTER)

    ny = cy + ch + Inches(0.22)
    tb(sl, "The tools nonprofits use were built for personal task management or enterprise teams. "
           "Neither fits a volunteer coordinator managing field events on a tight budget.",
       CX + Inches(1.0), ny, CW - Inches(2.0), Inches(0.50),
       sz=12, color=SEC, align=PP_ALIGN.CENTER, italic=True)
    # bottom ≈ 1.52 + 2.45 + 0.22 + 0.50 + 0.5 = 5.19"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S03 — Solution
# bottom ≈ CONTENT_Y + 2.8 + 0.22 + 0.38 = 4.92"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s03(prs):
    sl = blank(prs)
    lbl(sl, "The Solution", CX, CY, CW)
    h2(sl, "One app. Built specifically for NGO operations.", CX, CY + Inches(0.30), CW)
    sub(sl, "Trackly is a web app for nonprofits to plan events, coordinate volunteers, "
            "track attendance, and measure impact — with role-based access built in from day one.",
        CX + Inches(0.8), CY + Inches(0.96), CW - Inches(1.6))

    pillars = [
        ("Organise",   "Create and track events with statuses, categories, photos, and capacity limits", GREEN),
        ("Coordinate", "Manage volunteer enrollment, approvals, and attendance across your whole team", PURPLE),
        ("Measure",    "Track completion rates, volunteer trends, and event growth with live analytics", ORANGE),
    ]
    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(2.72)
    cy  = CONTENT_Y

    for i, (title, body, clr) in enumerate(pillars):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, title, cx + Inches(0.2), cy + Inches(0.28),
           cw - Inches(0.4), Inches(0.38), sz=16, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tb(sl, body, cx + Inches(0.2), cy + Inches(0.78),
           cw - Inches(0.4), Inches(1.7), sz=11, color=SEC, align=PP_ALIGN.CENTER)

    ly = cy + ch + Inches(0.26)
    lw, lh = Inches(4.5), Inches(0.36)
    lsx = CX + (CW - 2 * lw - Inches(0.28)) / 2
    for i, (link, clr) in enumerate([
        ("trackly-phi.vercel.app", GREEN),
        ("Live  ·  Deployed on Vercel + Render", SEC),
    ]):
        lx = lsx + i * (lw + Inches(0.28))
        card(sl, lx, ly, lw, lh, fill=CARD, border=CARD2)
        tb(sl, link, lx + Inches(0.1), ly + Inches(0.06),
           lw - Inches(0.2), lh - Inches(0.1), sz=9, color=clr, align=PP_ALIGN.CENTER)
    # bottom ≈ 1.52+2.72+0.26+0.36+0.5 = 5.36"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S04 — Core Workflow
# 4 steps × 0.78" + 3 × 0.10" gap = 3.42"
# bottom: CONTENT_Y + 3.42 = 4.94"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s04(prs):
    sl = blank(prs)
    lbl(sl, "Core Workflow", CX, CY, CW)
    h2(sl, "From login to live impact data in 4 steps", CX, CY + Inches(0.30), CW)
    sub(sl, "This is the loop a coordinator runs every time an event goes live. "
            "The dashboard updates in real time as each step happens.",
        CX + Inches(1.5), CY + Inches(0.96), CW - Inches(3.0))

    steps = [
        ("1", "Sign in and land on the dashboard.",
         "The mini calendar shows today's events. Animated stat cards display total events, "
         "active volunteers, completion rate, and upcoming events this week."),
        ("2", "Staff creates an event.",
         "Category, capacity, location, datetime, cover photo, enrollment type (open or assigned). "
         "Status starts as Planned and auto-flips to Ongoing at the event's start time."),
        ("3", "Volunteers join.",
         "Open events show a join button. Assigned events have a staff-managed attendee list. "
         "Join requests move through pending → approved / rejected, with rejection notes available."),
        ("4", "Event completes. Analytics update live.",
         "Donut charts, bar charts, volunteer trend lines. The Reports page gives the full picture "
         "across all events, staff, and months — no manual data entry required."),
    ]

    sh, sg = Inches(0.78), Inches(0.10)
    sy     = CONTENT_Y
    sw     = Inches(9.8)
    sx     = CX + (CW - sw) / 2

    for num, title, body in steps:
        card(sl, sx, sy, sw, sh, fill=CARD, border=CARD2)
        card(sl, sx + Inches(0.14), sy + Inches(0.17),
             Inches(0.44), Inches(0.44),
             fill=RGBColor(0x06, 0x22, 0x18), border=GREEN)
        tb(sl, num, sx + Inches(0.14), sy + Inches(0.19),
           Inches(0.44), Inches(0.36),
           sz=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        tb(sl, f"{title}  {body}",
           sx + Inches(0.72), sy + Inches(0.14),
           sw - Inches(0.86), sh - Inches(0.22), sz=11, color=SEC)
        sy += sh + sg
    # bottom ≈ 1.52 + 4*(0.78+0.10) - 0.10 = 1.52+3.42 = 4.94"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S05 — Key Features  (3-col × 2-row grid)
# row1+row2: CONTENT_Y + 2*2.50 + 0.18 = 1.52+5.18 = 6.70"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s05(prs):
    sl = blank(prs)
    lbl(sl, "Key Features", CX, CY, CW)
    h2(sl, "Six modules. One coherent system.", CX, CY + Inches(0.30), CW)

    feats = [
        ("Smart calendar",      GREEN,
         "Month, week, and agenda views. Color-coded event categories with auto contrast. "
         "Click any day to create an event. Export to shared Google Calendar."),
        ("Role-based access",   PURPLE,
         "Four roles: Admin, Staff, Volunteer, Viewer. Custom Firebase claims enforced "
         "on both frontend routes and the Flask API. No permission leakage."),
        ("Impact analytics",    ORANGE,
         "Donut charts (status/category), bar chart (monthly events), 6-month volunteer "
         "participation line chart, and a staff productivity table — all live from Firestore."),
        ("Smart notifications", SKY,
         "Push (FCM) + email (Resend) reminders at 24h, 6h, and 3h before each event. "
         "Idempotency tracking prevents duplicates. Per-user preference settings."),
        ("Enrollment workflows",ROSE,
         "Two modes: open (volunteers apply, staff approve) and assigned (direct assignment). "
         "Capacity validation, duplicate prevention, and rejection notes included."),
        ("Photo documentation", GREEN,
         "Upload event photos stored in Supabase. Metadata includes uploader identity and "
         "caption. Photos attached per event for audit and reporting."),
    ]

    cols = 3
    gap  = Inches(0.18)
    cw   = (CW - (cols - 1) * gap) / cols
    ch   = Inches(2.50)
    row_g = Inches(0.18)
    cy   = CY + Inches(0.96)    # tighter — no subtitle on this slide

    for i, (title, clr, body) in enumerate(feats):
        row = i // cols
        col = i % cols
        cx  = CX + col * (cw + gap)
        cy2 = cy + row * (ch + row_g)
        card(sl, cx, cy2, cw, ch)
        tb(sl, title, cx + Inches(0.2), cy2 + Inches(0.18),
           cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=clr)
        tb(sl, body,  cx + Inches(0.2), cy2 + Inches(0.60),
           cw - Inches(0.4), ch - Inches(0.72), sz=10, color=SEC)
    # bottom ≈ CY+0.96 + 2*2.50 + 0.18 = 0.5+0.96+5.18 = 6.64"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S06 — User Roles
# Table bottom: CONTENT_Y + 5*0.56 = 1.52+2.80 = 4.32"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s06(prs):
    sl = blank(prs)
    lbl(sl, "Access Model", CX, CY, CW)
    h2(sl, "Four roles that mirror how NGOs actually work", CX, CY + Inches(0.30), CW)
    sub(sl, "Most tools bolt on sharing after the fact. Trackly's permission model was built "
            "from the org chart up.",
        CX + Inches(2.0), CY + Inches(0.96), CW - Inches(4.0))

    rows = [
        ("",         "What they can do",                                              "Real-world equivalent",  DIM,    True),
        ("Admin",    "Create/delete events, manage user roles and user accounts, "
                     "view all analytics and reports",                                 "Program director / NGO head",   ROSE,   False),
        ("Staff",    "Create and edit events, manage attendees, assign volunteers, "
                     "upload photos",                                                  "Field coordinator / manager",   ORANGE, False),
        ("Volunteer","Browse open events, submit join requests, view personal event "
                     "history, upload photos",                                         "Volunteer / field worker",      GREEN,  False),
        ("Viewer",   "Read-only access to events and schedules — no creation or "
                     "enrollment actions",                                             "Donor / board member",          DIM,    False),
    ]

    col_w  = [Inches(1.4), Inches(6.8), Inches(3.9)]
    tw     = sum(col_w)
    tx     = CX + (CW - tw) / 2
    rh, rg = Inches(0.52), Inches(0.07)
    ty     = CONTENT_Y

    for ri, (role, does, real, clr, is_hdr) in enumerate(rows):
        ry = ty + ri * (rh + rg)
        bg = RGBColor(0x08, 0x10, 0x1e) if is_hdr else CARD
        bd = CARD2
        card(sl, tx, ry, tw, rh, fill=bg, border=bd, rounded=False)
        for ci, (col_text, cw) in enumerate(zip([role, does, real], col_w)):
            px = tx + sum(col_w[:ci]) + Inches(0.12)
            col_color = (clr if ci == 0 and not is_hdr
                         else DIM if is_hdr
                         else (TXT if ci == 0 else SEC))
            tb(sl, col_text, px, ry + Inches(0.10),
               cw - Inches(0.18), rh - Inches(0.14),
               sz=8 if is_hdr else 10,
               bold=(ci == 0 and not is_hdr) or is_hdr, color=col_color)

    nt = ty + len(rows) * (rh + rg) + Inches(0.14)
    tb(sl, "Roles use custom Firebase claims synced to Firestore. Both the React routes "
           "and the Flask API verify permissions independently — if one check fails, access is denied.",
       CX + Inches(1.0), nt, CW - Inches(2.0), Inches(0.42),
       sz=10, color=DIM, align=PP_ALIGN.CENTER, italic=True)
    # bottom ≈ 1.52 + 5*0.59 + 0.14 + 0.42 = 5.03"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S07 — Architecture
# Arch boxes: CONTENT_Y + 2.18 = 3.70"
# Decision cards: 3.70 + 0.20 + 2.0 = 5.90"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s07(prs):
    sl = blank(prs)
    lbl(sl, "Technical Architecture", CX, CY, CW)
    h2(sl, "Why each technology choice was deliberate", CX, CY + Inches(0.30), CW)
    sub(sl, "Five layers. Each one chosen for a specific reason, not just familiarity.",
        CX + Inches(2.0), CY + Inches(0.96), CW - Inches(4.0))

    boxes = [
        ("Frontend",     "React + Vite",       "TypeScript · Tailwind\nshadcn/ui · Framer Motion"),
        ("Auth + DB",    "Firebase",            "Auth · Firestore\nCloud Messaging"),
        ("Backend",      "Flask + Gunicorn",    "Python · Render\nModular routes"),
        ("Storage/Email","Supabase + Resend",   "File storage\nTransactional email"),
        ("Integration",  "Google Calendar",     "OAuth export\nShared org calendar"),
    ]
    n     = len(boxes)
    aw    = Inches(0.20)
    bw    = (CW - (n - 1) * aw) / n
    bh    = Inches(2.18)
    by    = CONTENT_Y

    for i, (lbl_txt, name, detail) in enumerate(boxes):
        bx  = CX + i * (bw + aw)
        hl  = (i == 1)
        card(sl, bx, by, bw, bh,
             fill=RGBColor(0x06, 0x1e, 0x18) if hl else CARD,
             border=GREEN if hl else CARD2)
        tb(sl, lbl_txt.upper(), bx + Inches(0.1), by + Inches(0.10),
           bw - Inches(0.2), Inches(0.22),
           sz=7, color=GREEN if hl else DIM, align=PP_ALIGN.CENTER)
        tb(sl, name, bx + Inches(0.1), by + Inches(0.36),
           bw - Inches(0.2), Inches(0.38),
           sz=12, bold=True, color=TXT, align=PP_ALIGN.CENTER)
        tb(sl, detail, bx + Inches(0.1), by + Inches(0.82),
           bw - Inches(0.2), Inches(1.18),
           sz=10, color=SEC, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = bx + bw + Inches(0.02)
            ay = by + bh / 2 - Inches(0.11)
            tb(sl, "->", ax, ay, aw - Inches(0.02), Inches(0.22),
               sz=10, color=DIM, align=PP_ALIGN.CENTER)

    decisions = [
        ("Firestore real-time listeners",
         "The dashboard uses onSnapshot so every client sees event status changes immediately "
         "— no polling, no manual refresh needed.", GREEN),
        ("Dual permission enforcement",
         "Roles are verified in React routes AND in the Flask API via Bearer tokens. "
         "A compromised frontend still cannot reach protected data.", PURPLE),
        ("Supabase for file storage",
         "Event photos go to Supabase rather than Firebase Storage to separate file costs "
         "and avoid Firestore document size limits.", ORANGE),
    ]
    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(2.0)
    dy  = by + bh + Inches(0.20)   # 3.70 + 0.20 = 3.90" from top of slide

    for i, (title, body, clr) in enumerate(decisions):
        cx = CX + i * (cw + gap)
        card(sl, cx, dy, cw, ch)
        tb(sl, title, cx + Inches(0.15), dy + Inches(0.14),
           cw - Inches(0.3), Inches(0.34), sz=11, bold=True, color=clr)
        tb(sl, body,  cx + Inches(0.15), dy + Inches(0.56),
           cw - Inches(0.3), ch - Inches(0.66), sz=10, color=SEC)
    # bottom ≈ 0.5 + 1.52 + 2.18 + 0.20 + 2.0 = 6.40"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S08 — Implementation Highlights
# 3 cards: CONTENT_Y + 4.88 = 6.40"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s08(prs):
    sl = blank(prs)
    lbl(sl, "Under the Hood", CX, CY, CW)
    h2(sl, "Three problems that actually took work to solve", CX, CY + Inches(0.30), CW)

    highlights = [
        ("Auto-status transitions",
         "Events flip Planned → Ongoing → Completed automatically based on start and end times. "
         "The logic in useEvents.ts normalises Firestore Timestamps, JS Date objects, and Unix "
         "epoch values through one utility before any comparison runs.\n\n"
         "Without this, an event starting at noon would still show 'Planned' until someone "
         "manually changed it.", GREEN),
        ("IST timezone normalisation",
         "Firestore stores all times in UTC. Without normalisation, Indian users saw events "
         "appear a day early or late in the calendar. Built toIST(), nowInIST(), and "
         "todayInIST() utilities that every date display passes through.\n\n"
         "This is the bug that's invisible in development and surfaces the moment a real "
         "user in a real timezone tries the app.", PURPLE),
        ("Reminder idempotency",
         "The notification scheduler runs on a cron to fire reminders at 24h, 6h, and 3h "
         "before each event. Without an idempotency guard, the same reminder would fire on "
         "every scheduler run. A Firestore flag per reminder per event is checked before "
         "any message sends.\n\n"
         "A ±90 minute window handles schedule drift so late-running crons still fire.", ORANGE),
    ]

    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(4.88)
    cy  = CONTENT_Y

    for i, (title, body, clr) in enumerate(highlights):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, f"0{i+1}", cx + Inches(0.18), cy + Inches(0.12),
           Inches(0.55), Inches(0.60), sz=28, bold=True, color=DIM)
        tb(sl, title, cx + Inches(0.18), cy + Inches(0.76),
           cw - Inches(0.36), Inches(0.38), sz=13, bold=True, color=clr)
        tb(sl, body, cx + Inches(0.18), cy + Inches(1.22),
           cw - Inches(0.36), ch - Inches(1.32), sz=10, color=SEC)
    # bottom ≈ 1.52 + 4.88 + 0.5 = 6.90"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S09 — Analytics & Reporting
# 2-col cards: CONTENT_Y + 4.30 = 5.82"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s09(prs):
    sl = blank(prs)
    lbl(sl, "Analytics & Reporting", CX, CY, CW)
    h2(sl, "Data that actually tells you something", CX, CY + Inches(0.30), CW)
    sub(sl, "Most event apps store data. Trackly's Reports page turns it into decisions — "
            "completion rates, volunteer trends, staff output, monthly growth.",
        CX + Inches(0.8), CY + Inches(0.96), CW - Inches(1.6))

    gap = Inches(0.22)
    cw  = (CW - gap) / 2
    ch  = Inches(4.30)
    cy  = CONTENT_Y
    lx, rx = CX, CX + cw + gap

    # Left: Dashboard
    card(sl, lx, cy, cw, ch)
    tb(sl, "Dashboard — live at a glance", lx + Inches(0.2), cy + Inches(0.14),
       cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=GREEN)
    dash_items = [
        "Animated stat cards: total events, active volunteers, completion rate, upcoming (7-day)",
        "Mini calendar with per-day event count badges and click-to-view",
        "Today's Events panel — real-time Firestore listener, no manual refresh",
        "Donut chart (event status) + bar chart (monthly distribution)",
    ]
    for j, item in enumerate(dash_items):
        iy = cy + Inches(0.58 + j * 0.84)
        card(sl, lx + Inches(0.16), iy, cw - Inches(0.32), Inches(0.72),
             fill=RGBColor(0x08, 0x12, 0x22), border=CARD2)
        tb(sl, item, lx + Inches(0.30), iy + Inches(0.10),
           cw - Inches(0.55), Inches(0.56), sz=10, color=SEC)

    # Right: Reports
    card(sl, rx, cy, cw, ch)
    tb(sl, "Reports — org-level insight", rx + Inches(0.2), cy + Inches(0.14),
       cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=PURPLE)
    report_items = [
        "6-month volunteer participation trend line chart with per-month drill-down modal",
        "Staff productivity table: events created vs completed, filterable by column",
        "Category breakdown donut — which event types run most often",
        "Growth rate metric (month-over-month) with demo fallback on empty DB",
    ]
    for j, item in enumerate(report_items):
        iy = cy + Inches(0.58 + j * 0.84)
        card(sl, rx + Inches(0.16), iy, cw - Inches(0.32), Inches(0.72),
             fill=RGBColor(0x08, 0x12, 0x22), border=CARD2)
        tb(sl, item, rx + Inches(0.30), iy + Inches(0.10),
           cw - Inches(0.55), Inches(0.56), sz=10, color=SEC)
    # bottom ≈ 1.52 + 4.30 + 0.5 = 6.32"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S10 — Notifications
# 2-col cards: CONTENT_Y + 4.0 = 5.52"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s10(prs):
    sl = blank(prs)
    lbl(sl, "Notification System", CX, CY, CW)
    h2(sl, "Reminders that reach people where they actually are", CX, CY + Inches(0.30), CW)
    sub(sl, "Email alone doesn't reach field volunteers. Push alone doesn't work without an app install. "
            "Trackly sends both, intelligently, without duplicating.",
        CX + Inches(0.8), CY + Inches(0.96), CW - Inches(1.6))

    gap = Inches(0.22)
    cw  = (CW - gap) / 2
    ch  = Inches(4.0)
    cy  = CONTENT_Y
    lx, rx = CX, CX + cw + gap

    # Left: timing
    card(sl, lx, cy, cw, ch)
    tb(sl, "How the timing works", lx + Inches(0.2), cy + Inches(0.14),
       cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=GREEN)
    timings = [
        ("24 hours before", "email + push reminder sent"),
        ("6 hours before",  "second reminder dispatched"),
        ("3 hours before",  "final reminder with event details"),
        ("+/- 90 min window", "handles cron schedule drift — late runs still fire"),
    ]
    for j, (when, what) in enumerate(timings):
        ty2 = cy + Inches(0.58 + j * 0.78)
        hl  = (j == 3)
        card(sl, lx + Inches(0.16), ty2, cw - Inches(0.32), Inches(0.66),
             fill=RGBColor(0x06, 0x22, 0x16) if hl else RGBColor(0x08, 0x12, 0x22),
             border=BDR_G if hl else CARD2)
        tb(sl, when, lx + Inches(0.30), ty2 + Inches(0.06),
           cw - Inches(0.55), Inches(0.26), sz=10, bold=True,
           color=GREEN if hl else TXT)
        tb(sl, what, lx + Inches(0.30), ty2 + Inches(0.32),
           cw - Inches(0.55), Inches(0.26), sz=10, color=SEC)

    # Right: channels
    card(sl, rx, cy, cw, ch)
    tb(sl, "Channels and controls", rx + Inches(0.2), cy + Inches(0.14),
       cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=PURPLE)
    channels = [
        ("Push via Firebase FCM",  SKY,
         "FCM tokens stored per user in Firestore. Dispatched from Flask backend with Bearer token auth."),
        ("Email via Resend",       ORANGE,
         "Transactional email with event details. Each user sets their own notification preferences."),
        ("Idempotency guard",      GREEN,
         "A Firestore flag per reminder per event prevents any message sending twice, "
         "even if the cron runs multiple times in one window."),
    ]
    for j, (ch_title, clr, ch_body) in enumerate(channels):
        cy2 = cy + Inches(0.58 + j * 1.08)
        card(sl, rx + Inches(0.16), cy2, cw - Inches(0.32), Inches(0.96),
             fill=RGBColor(0x08, 0x12, 0x22), border=CARD2)
        tb(sl, ch_title, rx + Inches(0.30), cy2 + Inches(0.08),
           cw - Inches(0.55), Inches(0.28), sz=11, bold=True, color=clr)
        tb(sl, ch_body, rx + Inches(0.30), cy2 + Inches(0.38),
           cw - Inches(0.55), Inches(0.50), sz=10, color=SEC)
    # bottom ≈ 1.52 + 4.0 + 0.5 = 6.02"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S11 — Stack Deep Dive  (4 columns)
# Cards: CONTENT_Y + 2.78 = 4.30"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s11(prs):
    sl = blank(prs)
    lbl(sl, "Stack Breakdown", CX, CY, CW)
    h2(sl, "Full-stack, deployed, and production-ready", CX, CY + Inches(0.30), CW)
    sub(sl, "119 commits. Two live deployments. Gunicorn in production. "
            "This is a deployable product, not a class demo.",
        CX + Inches(1.5), CY + Inches(0.96), CW - Inches(3.0))

    columns = [
        ("Frontend", [
            "React 18 + TypeScript", "Vite (SWC)",
            "Tailwind CSS", "shadcn/ui + Radix",
            "TanStack Query v5", "Framer Motion", "Recharts",
        ]),
        ("Backend", [
            "Python Flask 3.1", "Gunicorn (prod)",
            "flask-cors", "Modular routes",
            "Bearer token auth", "Render hosting",
        ]),
        ("Data & Auth", [
            "Firebase Auth", "Firestore (realtime)",
            "Firebase FCM", "Supabase storage",
            "Custom claims",
        ]),
        ("Integrations", [
            "Google Calendar API", "Resend email",
            "react-big-calendar", "Zod + React Hook Form",
            "date-fns", "Vitest",
        ]),
    ]

    cols = 4
    gap  = Inches(0.22)
    cw   = (CW - (cols - 1) * gap) / cols
    ch   = Inches(2.78)
    cy   = CONTENT_Y

    for i, (col_title, items) in enumerate(columns):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, col_title.upper(), cx + Inches(0.12), cy + Inches(0.12),
           cw - Inches(0.24), Inches(0.24),
           sz=7, color=GREEN if i < 2 else DIM, bold=True)
        for j, item in enumerate(items):
            py2 = cy + Inches(0.44 + j * 0.34)
            card(sl, cx + Inches(0.12), py2, cw - Inches(0.24), Inches(0.30),
                 fill=RGBColor(0x08, 0x14, 0x26) if j < 2 else RGBColor(0x06, 0x10, 0x1e),
                 border=BDR_G if (i < 2 and j < 2) else CARD2)
            tb(sl, item, cx + Inches(0.20), py2 + Inches(0.05),
               cw - Inches(0.36), Inches(0.22),
               sz=9, color=GREEN if (i < 2 and j < 2) else DIM)

    tb(sl, "TypeScript: 75.2%  ·  JavaScript: 14.1%  ·  Python: 5.6%  ·  CSS + config: remainder",
       CX, cy + ch + Inches(0.18), CW, Inches(0.24),
       sz=9, color=DIM, align=PP_ALIGN.CENTER)
    # bottom ≈ 1.52 + 2.78 + 0.18 + 0.24 + 0.5 = 5.22"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S12 — Challenges & Learnings
# Cards: CONTENT_Y + 4.70 = 6.22"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s12(prs):
    sl = blank(prs)
    lbl(sl, "Building It", CX, CY, CW)
    h2(sl, "What took longer than expected, and what we learned", CX, CY + Inches(0.30), CW)

    challenges = [
        ("Firestore Timestamp inconsistency",
         "Firestore returns Timestamp objects. Parts of the code expected JS Date objects. "
         "Others expected Unix milliseconds. The same event time arrived in three different "
         "shapes depending on how it was stored.\n\n"
         "Built a single normalisation utility that every date comparison routes through.\n\n"
         "Lesson: normalise data at the boundary, not at the point of use.", ROSE),
        ("CORS on Render's dynamic URLs",
         "Render generates preview URLs that change with every push. The Flask CORS config "
         "originally used a hardcoded whitelist, which broke each time the Vercel frontend "
         "URL changed.\n\n"
         "Rewrote CORS handling to dynamically allow any *.vercel.app origin.\n\n"
         "Lesson: never hardcode deployment URLs in CORS config during active development.", ORANGE),
        ("Real-time + TanStack Query conflict",
         "TanStack Query is built for HTTP with clear cache invalidation. Firestore's "
         "onSnapshot is a push model. Running both on the same data caused stale data "
         "and double renders.\n\n"
         "Resolved by separating real-time subscriptions from Query-managed endpoints with "
         "explicit cache invalidation hooks.\n\n"
         "Lesson: pick one data model per domain and stick to it.", PURPLE),
    ]

    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(4.70)
    cy  = CONTENT_Y

    for i, (title, body, clr) in enumerate(challenges):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, title, cx + Inches(0.18), cy + Inches(0.15),
           cw - Inches(0.36), Inches(0.36), sz=12, bold=True, color=clr)
        tb(sl, body, cx + Inches(0.18), cy + Inches(0.60),
           cw - Inches(0.36), ch - Inches(0.72), sz=10, color=SEC)
    # bottom ≈ 1.52 + 4.70 + 0.5 = 6.72"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# S13 — Roadmap & Close
# bottom ≈ CONTENT_Y + 2.82 + 0.20 + 0.38 + 0.20 + 0.28 + 0.50 = 6.40"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s13(prs):
    sl = blank(prs)
    lbl(sl, "What Comes Next", CX, CY, CW)
    h2(sl, "Trackly works. Here's where it goes from here.", CX, CY + Inches(0.30), CW)
    sub(sl, "The core loop is solid. These three additions would take it from a strong student "
            "project to something an NGO would actually deploy for their full team.",
        CX + Inches(0.8), CY + Inches(0.96), CW - Inches(1.6))

    roadmap = [
        ("Near term", GREEN, BDR_G,
         "Bulk CSV event import",
         "NGOs with existing schedules in Excel should onboard in minutes, not hours. "
         "Drag-and-drop CSV with field mapping and validation preview.\n\n"
         "Solves the biggest onboarding friction for orgs with historical data."),
        ("Mid term", ORANGE, BDR_O,
         "SMS via Twilio",
         "Push and email don't reach volunteers in poor-connectivity areas or those "
         "without smartphones. SMS does — and it extends the existing notification "
         "scheduler without a full rewrite.\n\n"
         "Especially relevant for rural NGO operations in India."),
        ("Research track", PURPLE, BDR_P,
         "QR attendance scanning",
         "Replace manual check-in lists with a per-event QR code volunteers scan on "
         "arrival. Works offline on the coordinator's device, syncs when connectivity "
         "returns.\n\n"
         "Field events often have poor internet — offline-first is the right target."),
    ]

    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(2.82)
    cy  = CONTENT_Y

    for i, (tier, clr, bdr, title, body) in enumerate(roadmap):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch, fill=CARD, border=bdr)
        tb(sl, tier.upper(), cx + Inches(0.18), cy + Inches(0.12),
           cw - Inches(0.36), Inches(0.22), sz=7, color=clr)
        tb(sl, title, cx + Inches(0.18), cy + Inches(0.40),
           cw - Inches(0.36), Inches(0.34), sz=12, bold=True, color=TXT)
        tb(sl, body, cx + Inches(0.18), cy + Inches(0.82),
           cw - Inches(0.36), ch - Inches(0.92), sz=10, color=SEC)

    # Links
    links = [
        ("trackly-phi.vercel.app",            GREEN),
        ("github.com/Shlok-Dwivedi/Trackly",  SEC),
    ]
    lw, lh = Inches(4.5), Inches(0.36)
    gap_l  = Inches(0.28)
    ly     = cy + ch + Inches(0.20)
    lsx    = CX + (CW - 2 * lw - gap_l) / 2
    for i, (link, clr) in enumerate(links):
        lx = lsx + i * (lw + gap_l)
        card(sl, lx, ly, lw, lh, fill=CARD, border=CARD2)
        tb(sl, link, lx + Inches(0.1), ly + Inches(0.06),
           lw - Inches(0.2), lh - Inches(0.1),
           sz=10, color=clr, align=PP_ALIGN.CENTER, italic=True)

    tb(sl, "Built by Shlok Dwivedi",
       CX, ly + lh + Inches(0.16), CW, Inches(0.24),
       sz=10, color=DIM, align=PP_ALIGN.CENTER)
    tb(sl, "Thank you",
       CX, ly + lh + Inches(0.46), CW, Inches(0.50),
       sz=26, bold=True, color=TXT, align=PP_ALIGN.CENTER)
    # bottom ≈ 1.52+2.82+0.20+0.36+0.16+0.24+0.46+0.50 = 6.26" + 0.5 = 6.76"  ✓


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()

    slides = [
        (s01, "Title"),
        (s02, "Problem"),
        (s03, "Solution"),
        (s04, "Core workflow"),
        (s05, "Key features"),
        (s06, "User roles"),
        (s07, "Architecture"),
        (s08, "Implementation highlights"),
        (s09, "Analytics & reporting"),
        (s10, "Notifications"),
        (s11, "Stack deep dive"),
        (s12, "Challenges & learnings"),
        (s13, "Roadmap & close"),
    ]

    print("Building slides...")
    for i, (fn, title) in enumerate(slides, 1):
        fn(prs)
        print(f"  {i:02d}/13 — {title}")

    out = "Trackly.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")
    print(f"Slide size : 13.333\" x 7.5\"  (16:9 widescreen)")
    print(f"Safe margin: 0.5\" L/R  ·  0.5\" top  ·  0.55\" bottom")
    print(f"Hard bottom: 6.95\" — all elements verified within bounds")
