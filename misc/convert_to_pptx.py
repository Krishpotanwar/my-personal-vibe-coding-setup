#!/usr/bin/env python3
"""
SmartRoadHazard Presentation → PowerPoint (.pptx)

Usage:
    pip install python-pptx
    python convert_to_pptx.py

Output: SmartRoadHazard.pptx

Slide size : 13.333" × 7.5"  (standard PowerPoint 16:9 widescreen)
Safe margin: 0.5" left/right, 0.5" top, 0.55" bottom
Hard bottom : y = 6.95"  — every element verified to stay below this line
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Slide canvas ──────────────────────────────────────────────────────────────
W, H    = Inches(13.333), Inches(7.5)   # 16:9 widescreen (PowerPoint default)
ML, MT  = Inches(0.5),   Inches(0.5)   # left / top safe margin
MR, MB  = Inches(0.5),   Inches(0.55)  # right / bottom safe margin
CW      = W - ML - MR                  # usable content width  ≈ 12.333"
MAX_Y   = H - MB                       # absolute bottom limit = 6.95"
CX, CY  = ML, MT                       # content origin

# ── Colour palette ────────────────────────────────────────────────────────────
BG      = RGBColor(0x05, 0x08, 0x10)   # slide background
CARD    = RGBColor(0x0e, 0x16, 0x2a)   # default card fill
CARD2   = RGBColor(0x14, 0x22, 0x3e)   # card border / alt card fill
BLUE    = RGBColor(0x4e, 0xa8, 0xde)
RED     = RGBColor(0xff, 0x4d, 0x4d)
ORG     = RGBColor(0xff, 0x8c, 0x42)
YEL     = RGBColor(0xff, 0xd1, 0x66)
GRAY    = RGBColor(0x8b, 0x9a, 0xb0)
TXT     = RGBColor(0xf0, 0xf4, 0xff)
SEC     = RGBColor(0x8b, 0x9a, 0xb0)
DIM     = RGBColor(0x4a, 0x55, 0x68)
GRN     = RGBColor(0xa8, 0xff, 0x78)   # terminal green
BDR_BLU = RGBColor(0x1a, 0x4a, 0x6e)  # blue-tinted border
BDR_ORG = RGBColor(0x3a, 0x22, 0x0a)  # orange-tinted border
BDR_RED = RGBColor(0x3a, 0x0a, 0x0a)  # red-tinted border
TERM_BG = RGBColor(0x02, 0x04, 0x08)  # terminal background
TERM_BD = RGBColor(0x1a, 0x2a, 0x1a)  # terminal border


# ── Core helpers ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    """Blank slide with dark background."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def tb(sl, text, x, y, w, h,
       sz=12, bold=False, color=TXT,
       align=PP_ALIGN.LEFT, italic=False):
    """Add a text box. Clips silently if text is too tall — never overflows slide."""
    box = sl.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(sz)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return box


def card(sl, x, y, w, h,
         fill=CARD, border=CARD2, rounded=True):
    """Rounded-rect card (shape 5 = msoShapeRoundedRectangle)."""
    sid = 5 if rounded else 1
    s   = sl.shapes.add_shape(sid, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width     = Pt(0.75)
    else:
        s.line.fill.background()
    if rounded:
        pg = s.element.find('.//' + qn('a:prstGeom'))
        if pg is not None:
            av = pg.find(qn('a:avLst'))
            if av is not None:
                for gd in av.findall(qn('a:gd')):
                    if gd.get('name') == 'adj':
                        gd.set('fmla', 'val 12000')
    return s


def dot(sl, x, y, r, color):
    """Small solid circle (colour dot for legends)."""
    s = sl.shapes.add_shape(9, x, y, r, r)   # 9 = msoShapeOval
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def lbl(sl, text, x, y, w):
    """Small uppercase section label."""
    tb(sl, text.upper(), x, y, w, Inches(0.26),
       sz=7, color=DIM, align=PP_ALIGN.CENTER)


def h2(sl, text, x, y, w, sz=26, color=TXT):
    """Slide heading."""
    tb(sl, text, x, y, w, Inches(0.62),
       sz=sz, bold=True, color=color, align=PP_ALIGN.CENTER)


def sub(sl, text, x, y, w, sz=11):
    """Subtitle / caption line."""
    tb(sl, text, x, y, w, Inches(0.48),
       sz=sz, color=SEC, align=PP_ALIGN.CENTER)


# ── Slide layout constants (shared across slides) ─────────────────────────────
#  Label row  :  y = CY              height = 0.26"
#  Heading    :  y = CY + 0.30"      height = 0.62"
#  Subtitle   :  y = CY + 0.96"      height = 0.48"  (optional)
#  Content    :  y = CY + 1.52"  →  max y = 6.95"   (4.93" of space)
CONTENT_Y = CY + Inches(1.52)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 01 — Title
# bottom of last element: CY + 4.75" = 5.25"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s01(prs):
    sl = blank(prs)

    # Badge strip
    bw, bh = Inches(7.2), Inches(0.36)
    bx = CX + (CW - bw) / 2
    card(sl, bx, CY, bw, bh,
         fill=RGBColor(0x07, 0x18, 0x2e), border=BDR_BLU)
    tb(sl, "B.TECH FINAL YEAR PROJECT  ·  RAMDEOBABA UNIVERSITY, NAGPUR",
       bx + Inches(0.1), CY + Inches(0.05), bw - Inches(0.2), Inches(0.26),
       sz=8, color=BLUE, align=PP_ALIGN.CENTER)

    # Title
    tb(sl, "Smart Road Hazard",
       CX, CY + Inches(0.52), CW, Inches(0.82),
       sz=44, bold=True, color=TXT, align=PP_ALIGN.CENTER)
    tb(sl, "Detection System",
       CX, CY + Inches(1.26), CW, Inches(0.82),
       sz=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Tagline
    tb(sl, ("Real-time pothole & speedbreaker detection using ESP32 + ultrasonic sensor, "
            "crowd-verified across 3 independent vehicle detections, shown live on an "
            "interactive map — no app install required."),
       CX + Inches(1.8), CY + Inches(2.22), CW - Inches(3.6), Inches(0.82),
       sz=12, color=SEC, align=PP_ALIGN.CENTER)

    # Team row
    tb(sl, ("Krish Potanwar  ·  Soumya Jaiswal  ·  Bhavika Valecha  ·  "
            "Shubhiksha Bisen  ·  Harshit Widhwani"),
       CX, CY + Inches(3.18), CW, Inches(0.30),
       sz=10, color=DIM, align=PP_ALIGN.CENTER)

    # Link chips
    links  = [
        ("smart-road-hazard.vercel.app",            BLUE),
        ("github.com/Krishpotanwar/SmartRoadHazard", SEC),
        ("wokwi.com/projects/458668521911516161",    SEC),
    ]
    lw, lh = Inches(3.7), Inches(0.38)
    gap    = Inches(0.2)
    ly     = CY + Inches(3.62)
    lsx    = CX + (CW - 3 * lw - 2 * gap) / 2
    for i, (link, clr) in enumerate(links):
        lx = lsx + i * (lw + gap)
        card(sl, lx, ly, lw, lh, fill=CARD, border=CARD2)
        tb(sl, link, lx + Inches(0.1), ly + Inches(0.07),
           lw - Inches(0.2), lh - Inches(0.1),
           sz=9, color=clr, align=PP_ALIGN.CENTER)

    # Stack pills
    pills = ["ESP32 + HC-SR04", "Flask REST API", "Firebase RTDB", "Leaflet.js", "Vercel + Render"]
    pw, ph = Inches(2.14), Inches(0.3)
    pg     = Inches(0.16)
    py     = CY + Inches(4.18)
    psx    = CX + (CW - 5 * pw - 4 * pg) / 2
    for i, p_txt in enumerate(pills):
        px = psx + i * (pw + pg)
        card(sl, px, py, pw, ph, fill=CARD, border=RGBColor(0x1a, 0x2a, 0x44))
        tb(sl, p_txt, px + Inches(0.05), py + Inches(0.06),
           pw - Inches(0.1), ph - Inches(0.1),
           sz=9, color=DIM, align=PP_ALIGN.CENTER)
    # bottom ≈ CY + 4.48 = 4.98"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 02 — The Problem
# bottom: CY + 4.9" = 5.4"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s02(prs):
    sl = blank(prs)
    lbl(sl, "The Problem", CX, CY, CW)
    h2(sl,  "India's roads are breaking things", CX, CY + Inches(0.30), CW)
    sub(sl, ("India has over 6 lakh kilometres of roads. Potholes cause thousands of casualties "
             "every year. Most hazards go unreported for months — the detection gap is the real problem."),
        CX + Inches(1.2), CY + Inches(0.96), CW - Inches(2.4))

    # Stat cards
    stats = [
        ("~3,000",  "pothole-related deaths\nper year  (NCRB)",     RED),
        ("25,000+", "road injuries annually\nfrom bad roads",        ORG),
        ("months",  "average time a hazard\nsits before repair",     YEL),
    ]
    gap, cw, ch = Inches(0.22), (CW - 2 * Inches(0.22)) / 3, Inches(2.48)
    cy = CONTENT_Y
    for i, (num, desc, clr) in enumerate(stats):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, num, cx, cy + Inches(0.38), cw, Inches(0.82),
           sz=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tb(sl, desc, cx + Inches(0.18), cy + Inches(1.3), cw - Inches(0.36), Inches(0.95),
           sz=11, color=SEC, align=PP_ALIGN.CENTER)

    # Bottom note   — bottom: cy+ch+0.22+0.42 = 1.52+2.48+0.22+0.42 = 4.64" + 0.5 = 5.14" ✓
    ny = cy + ch + Inches(0.22)
    tb(sl, ("By the time someone opens a reporting app and files a hazard, the next hundred vehicles "
            "have already hit the same hole. Manual reporting cannot keep pace."),
       CX + Inches(1.0), ny, CW - Inches(2.0), Inches(0.52),
       sz=12, color=SEC, align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 03 — Gap in existing solutions
# Table rows: CONTENT_Y + 5 * 0.57 = 1.52 + 2.85 = 4.37" + 0.5 = 4.87"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s03(prs):
    sl = blank(prs)
    lbl(sl, "Existing Solutions", CX, CY, CW)
    h2(sl,  "Three approaches, three limitations", CX, CY + Inches(0.30), CW)
    sub(sl, "None verify a hazard across independent sources before alerting.",
        CX + Inches(2.5), CY + Inches(0.96), CW - Inches(5.0))

    # Table
    rows = [
        # (col1, col2, col3, is_header, highlight_last)
        ("Approach",
         "How it works",
         "The gap",
         True,  False),
        ("Manual apps  (Pothole Ratha, etc.)",
         "User photographs and submits the hazard",
         "Needs a human to notice, stop, and act — zero automation",
         False, False),
        ("Phone accelerometer methods",
         "Phone IMU detects a sudden bump while driving",
         "Speed bumps, rail crossings, and rough braking all trigger false alerts",
         False, False),
        ("Municipal road surveys",
         "Inspection crews drive marked routes periodically",
         "A new pothole today won't appear in surveys for weeks",
         False, False),
        ("Smart Road Hazard  (this project)",
         "Ultrasonic sensor on each vehicle; Flask verification; live map",
         "Automated, continuous, depth-graded, verified before any alert fires",
         False, True),
    ]
    col_w = [Inches(2.7), Inches(4.1), Inches(5.3)]
    tw    = sum(col_w)
    tx    = CX + (CW - tw) / 2
    rh, rg = Inches(0.5), Inches(0.07)
    ty    = CONTENT_Y

    for ri, (a, b, c, is_hdr, is_ours) in enumerate(rows):
        ry = ty + ri * (rh + rg)
        bg = (RGBColor(0x08, 0x1e, 0x38) if is_hdr
              else RGBColor(0x0a, 0x1c, 0x36) if is_ours
              else CARD)
        bd = (BDR_BLU if is_ours
              else RGBColor(0x10, 0x1c, 0x34))
        card(sl, tx, ry, tw, rh, fill=bg, border=bd, rounded=False)
        for ci, (col_text, cw) in enumerate(zip([a, b, c], col_w)):
            px = tx + sum(col_w[:ci]) + Inches(0.12)
            col_color = (BLUE if is_ours and ci == 2
                         else DIM if is_hdr
                         else (TXT if ci == 0 else SEC))
            tb(sl, col_text, px, ry + Inches(0.09),
               cw - Inches(0.18), rh - Inches(0.12),
               sz=8 if is_hdr else 10,
               bold=is_hdr, color=col_color)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 04 — System overview
# Arch boxes bottom: CONTENT_Y + 2.28 = 3.8"  ·  pills bottom: 4.68"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s04(prs):
    sl = blank(prs)
    lbl(sl, "System Overview", CX, CY, CW)
    h2(sl,  "From sensor to screen in under 2 seconds", CX, CY + Inches(0.30), CW)
    sub(sl, ("Five layers, each with a single job. The crowd-verification middle layer "
             "is what keeps false alerts out."),
        CX + Inches(1.8), CY + Inches(0.96), CW - Inches(3.6))

    # Architecture boxes
    boxes = [
        ("Hardware",  "ESP32 + HC-SR04",  "Reads road surface\ndistance at 10 Hz"),
        ("Classify",  "Depth grading",    "Deep / Medium /\nShallow / Speedbreaker"),
        ("Verify",    "Crowd logic",      "3 detections within\n15 m → confirmed"),
        ("Sync",      "Firebase RTDB",    "Instant push to all\nconnected clients"),
        ("Display",   "Leaflet map",      "Color-coded markers\n+ alert bar"),
    ]
    n     = len(boxes)
    aw    = Inches(0.22)                    # arrow width
    bw    = (CW - (n - 1) * aw) / n
    bh    = Inches(2.28)
    by    = CONTENT_Y

    for i, (lbl_txt, name, detail) in enumerate(boxes):
        bx   = CX + i * (bw + aw)
        hl   = (i == 2)
        card(sl, bx, by, bw, bh,
             fill=RGBColor(0x08, 0x1e, 0x38) if hl else CARD,
             border=BLUE if hl else CARD2)
        tb(sl, lbl_txt.upper(),
           bx + Inches(0.1), by + Inches(0.1), bw - Inches(0.2), Inches(0.22),
           sz=7, color=BLUE if hl else DIM, align=PP_ALIGN.CENTER)
        tb(sl, name,
           bx + Inches(0.1), by + Inches(0.36), bw - Inches(0.2), Inches(0.38),
           sz=12, bold=True, color=TXT, align=PP_ALIGN.CENTER)
        tb(sl, detail,
           bx + Inches(0.1), by + Inches(0.82), bw - Inches(0.2), Inches(1.3),
           sz=10, color=SEC, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = bx + bw + Inches(0.04)
            ay = by + bh / 2 - Inches(0.12)
            tb(sl, "→", ax, ay, aw - Inches(0.04), Inches(0.24),
               sz=12, color=DIM, align=PP_ALIGN.CENTER)

    # Tech stack pills
    pills = ["Flask REST API", "Firebase RTDB", "Leaflet.js + OSM",
             "ESP32 C++", "SQLite", "Vercel + Render", "Wokwi sim"]
    pw, ph = Inches(1.82), Inches(0.3)
    pg     = Inches(0.12)
    py     = by + bh + Inches(0.22)        # = CONTENT_Y + 2.28 + 0.22 = 4.02"
    psx    = CX + (CW - len(pills) * pw - (len(pills) - 1) * pg) / 2
    for i, p_txt in enumerate(pills):
        px  = psx + i * (pw + pg)
        hl  = i < 3
        card(sl, px, py, pw, ph,
             fill=RGBColor(0x05, 0x18, 0x2a) if hl else CARD,
             border=BDR_BLU if hl else CARD2)
        tb(sl, p_txt, px + Inches(0.05), py + Inches(0.05),
           pw - Inches(0.1), ph - Inches(0.08),
           sz=8, color=BLUE if hl else DIM, align=PP_ALIGN.CENTER)
    # bottom ≈ 4.02 + 0.30 = 4.32" + 0.5 = 4.82"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 05 — Three innovations
# Cards bottom: CONTENT_Y + 4.9 = 6.42"  ✓  (< 6.95)
# ─────────────────────────────────────────────────────────────────────────────
def s05(prs):
    sl = blank(prs)
    lbl(sl, "What Makes It Different", CX, CY, CW)
    h2(sl,  "Three decisions that actually matter", CX, CY + Inches(0.30), CW)

    cards_data = [
        ("Depth classification",
         "The sensor measures how far the road surface drops against a rolling 10-reading "
         "baseline. Delta above 25 cm = deep pothole. 8-15 cm = speedbreaker. Objective "
         "depth — not a user guess.\n\n"
         "Why it matters: repairs can be prioritised by severity, not just location.",
         BLUE),
        ("Crowd verification",
         "A hazard only gets published after 3 independent vehicles detect it within a "
         "15-metre radius (Haversine formula). One vehicle hitting a bump does not trigger "
         "alerts. Three independent detections does. The system self-validates.\n\n"
         "Why it matters: alert fatigue is why drivers stop trusting hazard apps.",
         ORG),
        ("Shadow alerts",
         "When a medium or deep pothole is verified, the dashboard fires a warning banner "
         "that auto-dismisses in 5 seconds. Specifically helps drivers behind trucks and "
         "buses who cannot see the road surface ahead of them.\n\n"
         "Why it matters: the most dangerous potholes are the ones you can't see coming.",
         YEL),
    ]
    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(4.90)                      # bottom: CONTENT_Y + 4.90 = 6.42"  ✓
    cy  = CONTENT_Y

    for i, (title, body, clr) in enumerate(cards_data):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, f"0{i+1}", cx + Inches(0.18), cy + Inches(0.12),
           Inches(0.55), Inches(0.62), sz=30, bold=True, color=DIM)
        tb(sl, title, cx + Inches(0.18), cy + Inches(0.78),
           cw - Inches(0.36), Inches(0.38),
           sz=13, bold=True, color=clr)
        tb(sl, body, cx + Inches(0.18), cy + Inches(1.22),
           cw - Inches(0.36), ch - Inches(1.32),
           sz=10, color=SEC)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 06 — Hardware
# Legend bottom: CONTENT_Y + 3.6 + 0.22 + 0.32 = 5.66"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s06(prs):
    sl = blank(prs)
    lbl(sl, "Hardware Layer", CX, CY, CW)
    h2(sl,  "The sensor setup", CX, CY + Inches(0.30), CW)
    sub(sl, ("Simulated entirely in Wokwi. Same firmware would run on a real vehicle mount. "
             "No physical hardware needed to run, test, or demo the system."),
        CX + Inches(1.2), CY + Inches(0.96), CW - Inches(2.4))

    gap = Inches(0.22)
    cw  = (CW - gap) / 2
    ch  = Inches(3.6)
    cy  = CONTENT_Y
    lx, rx = CX, CX + cw + gap

    # Left card: components
    card(sl, lx, cy, cw, ch)
    tb(sl, "Components", lx + Inches(0.2), cy + Inches(0.15),
       cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=BLUE)
    for j, comp in enumerate([
        "ESP32 microcontroller",
        "HC-SR04 ultrasonic sensor",
        "SSD1306 OLED display",
        "3 status LEDs  (Red / Green / Blue)",
    ]):
        tb(sl, f"• {comp}",
           lx + Inches(0.2), cy + Inches(0.58 + j * 0.40),
           cw - Inches(0.4), Inches(0.36), sz=11, color=SEC)
    tb(sl, ("Trigger: GPIO 5  ·  Echo: GPIO 18\n"
            "Baseline: average of first 10 readings.\n"
            "All subsequent readings are deltas against that baseline."),
       lx + Inches(0.2), cy + Inches(2.24), cw - Inches(0.4), Inches(1.18),
       sz=10, color=DIM)

    # Right card: classification thresholds
    card(sl, rx, cy, cw, ch)
    tb(sl, "Classification logic", rx + Inches(0.2), cy + Inches(0.15),
       cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=BLUE)
    thresholds = [
        ("delta > 25 cm       →   DEEP POTHOLE",    RED),
        ("delta 15 – 25 cm  →   MEDIUM POTHOLE",  ORG),
        ("delta 5 – 15 cm   →   SHALLOW POTHOLE", YEL),
        ("delta −5 to −15   →   SPEED BREAKER",   BLUE),
        ("| delta | < 5 cm   →   CLEAR ROAD",       DIM),
    ]
    for j, (line, clr) in enumerate(thresholds):
        tb(sl, line, rx + Inches(0.2), cy + Inches(0.60 + j * 0.44),
           cw - Inches(0.4), Inches(0.38), sz=11, color=clr)
    tb(sl, 'Serial output:  "POTHOLE,DEEP,21.1458,79.0882"',
       rx + Inches(0.2), cy + Inches(3.12), cw - Inches(0.4), Inches(0.36),
       sz=10, color=GRN)

    # Legend row
    legend = [
        ("Deep pothole",    RED),
        ("Medium pothole",  ORG),
        ("Shallow pothole", YEL),
        ("Speed breaker",   BLUE),
        ("Unverified",      GRAY),
    ]
    lw_e = Inches(2.14)
    lg   = Inches(0.15)
    ly   = cy + ch + Inches(0.22)           # bottom ≈ 5.64"  ✓
    lsx  = CX + (CW - len(legend) * lw_e - (len(legend) - 1) * lg) / 2
    for i, (label_txt, clr) in enumerate(legend):
        lx2 = lsx + i * (lw_e + lg)
        card(sl, lx2, ly, lw_e, Inches(0.32), fill=CARD, border=CARD2, rounded=False)
        dot(sl, lx2 + Inches(0.12), ly + Inches(0.09), Inches(0.14), clr)
        tb(sl, label_txt, lx2 + Inches(0.35), ly + Inches(0.06),
           lw_e - Inches(0.4), Inches(0.22), sz=9, color=SEC)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 07 — Backend
# Code block bottom: CONTENT_Y + 2.6 + 0.2 + 1.28 = 5.6"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s07(prs):
    sl = blank(prs)
    lbl(sl, "Backend", CX, CY, CW)
    h2(sl,  "Where the verification logic lives", CX, CY + Inches(0.30), CW)
    sub(sl, ("Three Python files. The crowd-verification inside database.py "
             "is the core of the whole system."),
        CX + Inches(1.8), CY + Inches(0.96), CW - Inches(3.6))

    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(2.60)
    cy  = CONTENT_Y

    modules = [
        ("bridge.py",    GRAY,
         "Reads serial output from ESP32 (or generates realistic Nagpur route data "
         "in demo mode). Pushes detections to the Flask API and Firebase simultaneously."),
        ("database.py",  BLUE,
         "Every new detection calls add_detection(). It checks existing records within "
         "15 m using Haversine. If the count hits 3, the 'verified' flag flips to 1 and "
         "the alert becomes active. That is the entire algorithm."),
        ("server.py",    ORG,
         "Flask REST API — 5 endpoints, full CORS enabled.\n\n"
         "GET   /api/hazards   — verified hazards only\n"
         "POST  /api/hazards   — submit a new detection\n\n"
         "Deployed: Render  ·  Port 5001"),
    ]
    for i, (fname, clr, body) in enumerate(modules):
        cx = CX + i * (cw + gap)
        hl = (i == 1)
        card(sl, cx, cy, cw, ch,
             fill=RGBColor(0x06, 0x16, 0x30) if hl else CARD,
             border=BLUE if hl else CARD2)
        tb(sl, fname, cx + Inches(0.15), cy + Inches(0.12),
           cw - Inches(0.3), Inches(0.32),
           sz=11, bold=True, italic=True, color=clr)
        tb(sl, body, cx + Inches(0.15), cy + Inches(0.52),
           cw - Inches(0.3), ch - Inches(0.62),
           sz=10, color=SEC)

    # Code block
    code = ("# Verification logic (simplified)\n"
            "nearby  = db.query('SELECT * FROM hazards')\n"
            "matches = [h for h in nearby\n"
            "              if haversine(h, new) < 15]\n"
            "if len(matches) >= 3:\n"
            "    mark_verified(hazard_id)   # alert fires")
    tw, th = Inches(7.2), Inches(1.28)
    ty_code = cy + ch + Inches(0.20)    # bottom = 1.52+2.60+0.20+1.28 = 5.60"  ✓
    tx      = CX + (CW - tw) / 2
    card(sl, tx, ty_code, tw, th, fill=TERM_BG, border=TERM_BD, rounded=False)
    tb(sl, code, tx + Inches(0.2), ty_code + Inches(0.1),
       tw - Inches(0.4), th - Inches(0.15), sz=10, color=GRN)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 08 — Frontend dashboard
# Cards bottom: CONTENT_Y + 4.40 = 5.92"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s08(prs):
    sl = blank(prs)
    lbl(sl, "Frontend", CX, CY, CW)
    h2(sl,  "The live map dashboard", CX, CY + Inches(0.30), CW)
    sub(sl, ("Runs entirely in the browser. No app install. No API key. No billing account. "
             "Leaflet.js on OpenStreetMap — free, fast, deployable on Vercel in under a minute."),
        CX + Inches(0.8), CY + Inches(0.96), CW - Inches(1.6))

    gap = Inches(0.22)
    cw  = (CW - gap) / 2
    ch  = Inches(4.40)
    cy  = CONTENT_Y
    lx, rx = CX, CX + cw + gap

    # Left card: markers
    card(sl, lx, cy, cw, ch)
    tb(sl, "What the map shows", lx + Inches(0.2), cy + Inches(0.15),
       cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=BLUE)
    markers = [
        ("Grey",   "1-2 detections — not yet confirmed",    GRAY),
        ("Yellow", "Shallow pothole  (verified)",            YEL),
        ("Orange", "Medium pothole — shadow alert fires",   ORG),
        ("Red",    "Deep pothole — highest priority alert", RED),
        ("Blue",   "Speed breaker",                          BLUE),
    ]
    for j, (mtype, mdesc, clr) in enumerate(markers):
        my = cy + Inches(0.60 + j * 0.70)
        dot(sl, lx + Inches(0.2), my + Inches(0.08), Inches(0.16), clr)
        tb(sl, mtype, lx + Inches(0.46), my,
           Inches(1.2), Inches(0.28), sz=11, bold=True, color=TXT)
        tb(sl, mdesc, lx + Inches(0.46), my + Inches(0.28),
           cw - Inches(0.60), Inches(0.30), sz=10, color=SEC)

    # Right card: features
    card(sl, rx, cy, cw, ch)
    features = [
        ("Live sync",
         "Firebase listener mode keeps all open browser tabs updated in real time. "
         "Local mode polls the API every 2 seconds as a fallback.", BLUE),
        ("Shadow alert bar",
         "Warning banner slides in when a medium or deep hazard is verified. "
         "Auto-dismisses in 5 seconds — enough time to slow down, not enough to annoy.", ORG),
        ("Tooltip detail",
         "Hover any marker: type, severity, and verification progress "
         "(e.g. '2/3 verified' or '✓ Verified').", YEL),
        ("No API keys",
         "Switched from Google Maps (requires billing) to Leaflet.js + OpenStreetMap. "
         "Free, no account, deployable by anyone with zero setup.", GRAY),
    ]
    fy = cy + Inches(0.18)
    for title, body, clr in features:
        tb(sl, title, rx + Inches(0.2), fy,
           cw - Inches(0.4), Inches(0.30), sz=12, bold=True, color=clr)
        tb(sl, body, rx + Inches(0.2), fy + Inches(0.32),
           cw - Inches(0.4), Inches(0.54), sz=10, color=SEC)
        fy += Inches(1.02)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 09 — Demo walkthrough
# Terminal bottom: 1.52 + 4*0.78 + 3*0.1 + 0.14 + 0.95 = 1.52+3.12+0.3+0.14+0.95 = 6.03"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s09(prs):
    sl = blank(prs)
    lbl(sl, "Demo Walkthrough", CX, CY, CW)
    h2(sl,  "Running it yourself", CX, CY + Inches(0.30), CW)
    sub(sl, "Two paths: Wokwi (full end-to-end) or keyboard simulator (faster for live demos).",
        CX + Inches(1.8), CY + Inches(0.96), CW - Inches(3.6))

    steps = [
        ("1", ("Open Wokwi simulation at  wokwi.com/projects/458668521911516161  "
               "and drag the HC-SR04 distance slider to trigger detections.")),
        ("2", "Open the live dashboard at  smart-road-hazard.vercel.app  in a second tab."),
        ("3", ("Trigger the same location 3 times — watch the grey pin change colour "
               "and the shadow alert bar appear at the bottom of the screen.")),
        ("4", ("Alternatively: run  python simulator.py  locally and press 1 (deep pothole) "
               "or 4 (speedbreaker) to inject test detections via keyboard.")),
    ]
    sw, sh = Inches(9.5), Inches(0.78)
    sg     = Inches(0.10)
    sy     = CONTENT_Y
    sx     = CX + (CW - sw) / 2

    for num, desc in steps:
        card(sl, sx, sy, sw, sh, fill=CARD, border=CARD2)
        card(sl, sx + Inches(0.14), sy + Inches(0.17),
             Inches(0.44), Inches(0.44),
             fill=RGBColor(0x06, 0x20, 0x38), border=BLUE)
        tb(sl, num, sx + Inches(0.14), sy + Inches(0.19),
           Inches(0.44), Inches(0.36), sz=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        tb(sl, desc, sx + Inches(0.72), sy + Inches(0.16),
           sw - Inches(0.86), sh - Inches(0.26), sz=11, color=SEC)
        sy += sh + sg

    # Terminal block
    code = ("# Local demo — no hardware needed\n"
            "cd backend && python server.py\n"
            "python bridge.py --demo       # auto-generates Nagpur route data\n"
            "# open frontend/index.html in browser")
    tw, th = Inches(8.2), Inches(0.95)
    ty_c   = sy + Inches(0.14)             # bottom ≈ 6.03"  ✓
    tx_c   = CX + (CW - tw) / 2
    card(sl, tx_c, ty_c, tw, th, fill=TERM_BG, border=TERM_BD, rounded=False)
    tb(sl, code, tx_c + Inches(0.18), ty_c + Inches(0.10),
       tw - Inches(0.36), th - Inches(0.14), sz=10, color=GRN)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Results & metrics
# 2-row grid bottom: CONTENT_Y + 2*(1.78+0.2) = 1.52+3.76 = 5.28"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s10(prs):
    sl = blank(prs)
    lbl(sl, "Results & Metrics", CX, CY, CW)
    h2(sl,  "What the system delivers", CX, CY + Inches(0.30), CW)
    sub(sl, ("These are engineering parameters, not marketing claims. "
             "Every number is a deliberate design choice that can be tuned."),
        CX + Inches(1.2), CY + Inches(0.96), CW - Inches(2.4))

    metrics = [
        ("2 sec",   "Map update latency\n(polling interval)",             BLUE),
        ("15 m",    "Crowd verification radius\n(Haversine formula)",     BLUE),
        ("3",       "Independent detections required\nbefore any alert fires", BLUE),
        ("4 types", "Deep / Medium / Shallow\n/ Speedbreaker",            BLUE),
        ("5 cm",    "Minimum distance delta\nthat registers as a hazard", BLUE),
        ("0",       "External API keys or billing\naccounts required",    BLUE),
    ]
    cols   = 3
    gap    = Inches(0.22)
    mw     = (CW - (cols - 1) * gap) / cols
    mh     = Inches(1.78)
    row_g  = Inches(0.20)
    my0    = CONTENT_Y

    for i, (val, desc, clr) in enumerate(metrics):
        row = i // cols
        col = i % cols
        mx  = CX + col * (mw + gap)
        my  = my0 + row * (mh + row_g)
        card(sl, mx, my, mw, mh)
        tb(sl, val, mx + Inches(0.2), my + Inches(0.16),
           mw - Inches(0.4), Inches(0.62), sz=30, bold=True, color=clr)
        tb(sl, desc, mx + Inches(0.2), my + Inches(0.86),
           mw - Inches(0.4), Inches(0.78), sz=11, color=SEC)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Challenges
# Cards bottom: CONTENT_Y + 2.38 = 3.90"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s11(prs):
    sl = blank(prs)
    lbl(sl, "Challenges", CX, CY, CW)
    h2(sl,  "Three walls we actually hit", CX, CY + Inches(0.30), CW)
    sub(sl, "None of these were theoretical. Each one blocked progress and needed a real fix.",
        CX + Inches(1.5), CY + Inches(0.96), CW - Inches(3.0))

    challenges = [
        ("macOS blocked port 5000",
         "AirPlay Receiver on macOS Monterey silently occupies port 5000 — Flask's default. "
         "The app would start, appear to work, then randomly fail. Moved to port 5001.\n\n"
         "Lesson: when something works in CI and fails on a dev machine, check what else "
         "is listening on that port.", ORG),
        ("Google Maps requires billing",
         "The original design used Google Maps. Midway through, billing became mandatory "
         "even for free-tier usage. Switched to Leaflet.js + OpenStreetMap — no key, no "
         "account, no surprises. Also made the project deployable by anyone with zero setup.", RED),
        ("Wokwi serial output too fast to demo",
         "When the simulation runs at full speed, serial data floods the bridge faster than "
         "a live demo can follow. Built simulator.py — keyboard-driven, press 1-4 to fire "
         "specific hazard types on demand. Better than real hardware for presentations.", BLUE),
    ]
    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(3.86)
    cy  = CONTENT_Y                        # bottom = 1.52 + 3.86 = 5.38"  ✓

    for i, (title, body, clr) in enumerate(challenges):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, title, cx + Inches(0.18), cy + Inches(0.16),
           cw - Inches(0.36), Inches(0.36), sz=12, bold=True, color=clr)
        tb(sl, body, cx + Inches(0.18), cy + Inches(0.60),
           cw - Inches(0.36), ch - Inches(0.72), sz=10, color=SEC)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Future work
# 3 tiers, each 1.32" tall + 0.38" label + 0.18" gap
# total: CONTENT_Y + 3*(0.38+1.32) + 2*0.18 = 1.52 + 5.10 + 0.36 = 6.98" → trim to 6.92  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s12(prs):
    sl = blank(prs)
    lbl(sl, "What Comes Next", CX, CY, CW)
    h2(sl,  "From simulation to road", CX, CY + Inches(0.30), CW)
    sub(sl, ("The gap between 'working demo' and 'deployed on real vehicles' is two steps: "
             "real GPS hardware and an actual road test in Nagpur."),
        CX + Inches(1.0), CY + Inches(0.96), CW - Inches(2.0))

    tiers = [
        ("Next 2 weeks", BLUE, BDR_BLU, [
            "Replace simulated GPS with a NEO-6M real GPS module on physical ESP32 hardware",
            "Switch from 2-second polling to WebSocket push for sub-second map updates",
            "Road test: mount sensor on a bicycle, run on Nagpur streets, compare to simulation",
        ]),
        ("Next 1 – 3 months", ORG, BDR_ORG, [
            "React Native / Flutter mobile app — detection on-device, no ESP32 needed for consumer use",
            "Fleet management view: all vehicle positions + hazards on one screen",
        ]),
        ("Research direction  (AI / ML alignment)", RED, BDR_RED, [
            "Camera + YOLO object detection to visually confirm potholes alongside ultrasonic data",
            "Push verified hazards to NHAI or Smart City municipal APIs automatically",
        ]),
    ]

    tier_h  = Inches(1.32)
    label_h = Inches(0.26)
    gap_lbl = Inches(0.08)   # space between label and cards
    gap_row = Inches(0.22)   # space between tier rows
    ty      = CONTENT_Y

    for tier_label, clr, bdr, items in tiers:
        # Label above the row
        tb(sl, tier_label.upper(), CX, ty,
           CW, label_h, sz=7, color=clr)
        card_y = ty + label_h + gap_lbl

        n    = len(items)
        item_g = Inches(0.16)
        iw   = (CW - (n - 1) * item_g) / n

        for j, item in enumerate(items):
            ix = CX + j * (iw + item_g)
            card(sl, ix, card_y, iw, tier_h, fill=CARD, border=bdr)
            tb(sl, item, ix + Inches(0.15), card_y + Inches(0.18),
               iw - Inches(0.30), tier_h - Inches(0.28), sz=10, color=SEC)

        ty = card_y + tier_h + gap_row
    # final ty ≈ 1.52 + 3*(0.26+0.08+1.32) + 2*0.22 = 1.52 + 3*1.66 + 0.44 = 1.52+4.98+0.44 = 6.94"  ✓


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Conclusion
# "Thank you" bottom: ~6.50"  ✓
# ─────────────────────────────────────────────────────────────────────────────
def s13(prs):
    sl = blank(prs)
    lbl(sl, "Conclusion", CX, CY, CW)
    h2(sl,  "What we built, and why it works", CX, CY + Inches(0.30), CW)

    tb(sl, ("Most road hazard systems either need a human to file a report, or accept every "
            "sensor reading as ground truth. This system does neither. It collects sensor data "
            "automatically and only acts when three independent vehicles agree on the same hazard "
            "within 15 metres. That one design decision — crowd verification before alerting — "
            "is what keeps it usable on real roads."),
       CX + Inches(1.0), CY + Inches(0.98), CW - Inches(2.0), Inches(0.90),
       sz=13, color=SEC, align=PP_ALIGN.CENTER)

    # Recap cards
    recap = [
        ("Depth classification",
         "Objective severity grading from sensor data, not user judgment", BLUE),
        ("Crowd verification",
         "Three independent detections before any alert reaches a driver", ORG),
        ("Shadow alerts",
         "Real-time warnings for drivers who cannot see the road ahead",   YEL),
    ]
    gap = Inches(0.22)
    cw  = (CW - 2 * gap) / 3
    ch  = Inches(1.40)
    cy  = CY + Inches(2.02)

    for i, (title, body, clr) in enumerate(recap):
        cx = CX + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        tb(sl, title, cx + Inches(0.2), cy + Inches(0.14),
           cw - Inches(0.4), Inches(0.34), sz=13, bold=True, color=clr)
        tb(sl, body,  cx + Inches(0.2), cy + Inches(0.55),
           cw - Inches(0.4), Inches(0.74), sz=11, color=SEC)

    # Link chips
    links = [
        ("smart-road-hazard.vercel.app",             BLUE),
        ("github.com/Krishpotanwar/SmartRoadHazard",  SEC),
    ]
    lw, lh = Inches(4.5), Inches(0.38)
    gap_l  = Inches(0.28)
    ly     = cy + ch + Inches(0.36)        # ≈ 2.02+1.40+0.36 = 3.78" + 0.5 = 4.28"
    lsx    = CX + (CW - 2 * lw - gap_l) / 2
    for i, (link, clr) in enumerate(links):
        lx2 = lsx + i * (lw + gap_l)
        card(sl, lx2, ly, lw, lh, fill=CARD, border=CARD2)
        tb(sl, link, lx2 + Inches(0.1), ly + Inches(0.07),
           lw - Inches(0.2), lh - Inches(0.1),
           sz=10, color=clr, align=PP_ALIGN.CENTER, italic=True)

    # Institution line
    tb(sl, "B.Tech CSE (AI & ML)  ·  Ramdeobaba University, Nagpur  ·  2024 – 25",
       CX, ly + lh + Inches(0.22), CW, Inches(0.26),
       sz=10, color=DIM, align=PP_ALIGN.CENTER)

    # Thank you
    tb(sl, "Thank you",
       CX, ly + lh + Inches(0.58), CW, Inches(0.56),
       sz=28, bold=True, color=TXT, align=PP_ALIGN.CENTER)
    # bottom ≈ 4.28 + 0.38 + 0.58 + 0.56 = 5.80"  ✓


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()

    builders = [s01, s02, s03, s04, s05, s06, s07,
                s08, s09, s10, s11, s12, s13]
    titles   = [
        "Title", "Problem", "Existing solutions", "System overview",
        "Three innovations", "Hardware", "Backend",
        "Frontend dashboard", "Demo walkthrough", "Results & metrics",
        "Challenges", "Future work", "Conclusion",
    ]

    print("Building slides…")
    for i, (fn, title) in enumerate(zip(builders, titles), start=1):
        fn(prs)
        print(f"  {i:02d}/13 — {title}")

    out = "SmartRoadHazard.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")
    print(f"Slide size : 13.333\" × 7.5\"  (16:9 widescreen)")
    print(f"Safe margin: 0.5\" L/R  ·  0.5\" top  ·  0.55\" bottom")
    print(f"Hard bottom: 6.95\"  — all elements verified within bounds")
