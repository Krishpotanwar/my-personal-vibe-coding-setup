#!/usr/bin/env python3
"""
cv_to_pptx.py  —  Connected Vehicles Case Studies → PowerPoint
══════════════════════════════════════════════════════════════════
Generates  connected_vehicles.pptx  (16:9 widescreen, B&W theme)
with PowerPoint Morph transitions on every slide transition.

Install:  pip install python-pptx
Run:      python cv_to_pptx.py
Output:   connected_vehicles.pptx

Morph notes
───────────
  • Shapes prefixed with !! (e.g. !!watermark, !!accent-line, !!dot)
    are matched by name across slides so PowerPoint Morph animates
    them from position A → position B between slides.
  • Open in PowerPoint 2019 / Microsoft 365 — Morph only plays there.
  • LibreOffice Impress does NOT support Morph; slides will still open.
"""

from __future__ import annotations
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ══════════════════════════════════════════════════════════════
#  PALETTE
# ══════════════════════════════════════════════════════════════
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
G1    = RGBColor(0xF7, 0xF7, 0xF7)   # light card background
G2    = RGBColor(0xE4, 0xE4, 0xE4)   # borders / dividers
G3    = RGBColor(0xAA, 0xAA, 0xAA)   # muted / labels
G4    = RGBColor(0x55, 0x55, 0x55)   # section headers
G5    = RGBColor(0x22, 0x22, 0x22)   # body text
WM    = RGBColor(0xED, 0xED, 0xED)   # watermark (very light)

# ══════════════════════════════════════════════════════════════
#  CANVAS  —  16:9 widescreen  (13.333" × 7.5")
# ══════════════════════════════════════════════════════════════
SW  = Inches(13.333)   # slide width
SH  = Inches(7.500)    # slide height
ML  = Inches(0.70)     # left  margin
MT  = Inches(0.60)     # top   margin
MR  = Inches(0.70)     # right margin
MB  = Inches(0.45)     # bottom margin
CW  = SW - ML - MR    # content width  ≈ 11.93"
CH  = SH - MT - MB    # content height ≈ 6.45"

FONT = "Arial"         # swap to 'Helvetica Neue' if installed

# XML namespaces used for Morph injection
_PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_DNS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ══════════════════════════════════════════════════════════════
#  PRIMITIVES
# ══════════════════════════════════════════════════════════════

def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs: Presentation):
    sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl


def add_morph_transition(slide, dur_ms: int = 600) -> None:
    """
    Injects  <p:transition dur="600" advClick="1"><p:morph/></p:transition>
    into the slide XML.  Must be called AFTER all shapes are added.
    """
    tr = etree.SubElement(
        slide.element,
        "{%s}transition" % _PNS,
        attrib={"dur": str(dur_ms), "advClick": "1"},
    )
    etree.SubElement(tr, "{%s}morph" % _PNS)


# ── text box ──────────────────────────────────────────────────

def tb(
    slide, text: str,
    l, t, w, h,
    size: float = 11,
    bold: bool = False,
    color: RGBColor = BLACK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    italic: bool = False,
    name: str | None = None,
) -> object:
    box = slide.shapes.add_textbox(l, t, w, h)
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = FONT
    return box


# ── solid rectangle ───────────────────────────────────────────

def rect(
    slide, l, t, w, h,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    lw: Emu = Pt(0.5),
    name: str | None = None,
) -> object:
    shp = slide.shapes.add_shape(1, l, t, w, h)   # 1 = RECTANGLE
    if name:
        shp.name = name
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width     = lw
    else:
        shp.line.fill.background()
    return shp


def hline(slide, l, t, w, color: RGBColor = G2, name: str | None = None) -> object:
    """Thin horizontal rule (2 pt high)."""
    return rect(slide, l, t, w, Emu(int(Pt(2))), fill=color, name=name)


def vline(slide, l, t, h, color: RGBColor = G2, name: str | None = None) -> object:
    """Thin vertical rule (1.5 pt wide)."""
    return rect(slide, l, t, Emu(int(Pt(1.5))), h, fill=color, name=name)


# ══════════════════════════════════════════════════════════════
#  SLIDE FURNITURE  (reused on every slide)
# ══════════════════════════════════════════════════════════════

# Watermark positions per slide — same shape name = Morph animates it
_WM_POS: dict[str, tuple] = {
    "01": (Inches(7.20),  Inches(0.40),  250),
    "02": (Inches(0.00),  Inches(-0.20), 185),
    "03": (Inches(3.70),  Inches(0.20),  280),
    "04": (Inches(7.90),  Inches(1.90),  210),
    "05": (Inches(0.00),  Inches(0.80),  225),
    "06": (Inches(2.90),  Inches(1.30),  265),
}

# Accent-line position/size per slide (l, t, w, h)
_LINE_CFG: dict[str, tuple] = {
    "01": (ML,                SH - Inches(0.55), CW,           Emu(int(Pt(2.5)))),
    "02": (ML,                MT,               Emu(int(Pt(2.5))), CH          ),
    "03": (SW * 0.49,         MT + Inches(0.9), Emu(int(Pt(2.5))), CH - Inches(0.9)),
    "04": (ML,                MT,               CW * 0.44,    Emu(int(Pt(2.5)))),
    "05": (SW - Inches(0.35), MT,               Emu(int(Pt(2.5))), CH          ),
    "06": (ML,                SH - Inches(0.80), CW * 0.54,   Emu(int(Pt(2.5)))),
}

# Dot positions per slide
_DOT_POS: dict[str, tuple] = {
    "01": (ML,                MT               ),
    "02": (SW - Inches(1.0),  SH - Inches(1.0) ),
    "03": (ML + Inches(0.05), SH * 0.50        ),
    "04": (SW * 0.50,         MT + Inches(0.05)),
    "05": (SW * 0.50,         SH - Inches(0.75)),
    "06": (SW * 0.68,         SH - Inches(0.80)),
}


def add_morph_furniture(slide, slide_num: str) -> None:
    """
    Adds the three persistent morph elements to a slide.
    Their !! names cause PowerPoint to morph them between slides.
    """
    # 1. Watermark number
    l, t, fs = _WM_POS[slide_num]
    tb(slide, slide_num, l, t, Inches(6), Inches(5),
       size=fs, bold=True, color=WM, name="!!watermark")

    # 2. Accent line
    l, t, w, h = _LINE_CFG[slide_num]
    rect(slide, l, t, w, h, fill=BLACK, name="!!accent-line")

    # 3. Dot
    l, t = _DOT_POS[slide_num]
    dot_sz = Inches(0.09)
    shp = rect(slide, l, t, dot_sz, dot_sz, fill=BLACK, name="!!dot")
    # Make it a circle
    spPr = shp.element.find(".//{%s}spPr" % _DNS)
    if spPr is not None:
        pg = spPr.find("{%s}prstGeom" % _DNS)
        if pg is not None:
            pg.set("prst", "ellipse")


# ══════════════════════════════════════════════════════════════
#  REUSABLE WIDGETS
# ══════════════════════════════════════════════════════════════

def eyebrow(slide, text: str) -> None:
    tb(slide, text, ML, MT, CW, Inches(0.22), size=8, color=G3)


def heading(slide, text: str, size: float = 34) -> None:
    tb(slide, text, ML, MT + Inches(0.24), CW, Inches(0.65),
       size=size, bold=True)


def note_box(slide, text: str, l, t, w, h) -> None:
    """Left-bordered callout."""
    bar_w = Emu(int(Pt(3)))
    rect(slide, l, t, bar_w, h, fill=BLACK)
    rect(slide, l + bar_w, t, w - bar_w, h, fill=G1, line=G2)
    tb(slide, text,
       l + Inches(0.14), t + Inches(0.08),
       w - Inches(0.20), h - Inches(0.14),
       size=10, color=G5)


def stat_block(slide, num: str, label: str, l, t, w) -> None:
    """Big number + small label inside a bordered box."""
    h = Inches(0.85)
    rect(slide, l, t, w, h, line=G2)
    tb(slide, num,   l + Inches(0.1), t + Inches(0.06), w - Inches(0.15),
       Inches(0.48), size=28, bold=True)
    tb(slide, label, l + Inches(0.1), t + Inches(0.54), w - Inches(0.15),
       Inches(0.28), size=8, color=G3)


def bullet_list(slide, items: list[str], l, t, w,
                row_h = Inches(0.40), size: float = 10) -> None:
    for i, item in enumerate(items):
        y = t + i * row_h
        tb(slide, "—",   l,                y, Inches(0.17), row_h, size=size, color=G3)
        tb(slide, item,  l + Inches(0.19), y, w - Inches(0.19), row_h,
           size=size, color=G5)
        if i < len(items) - 1:
            hline(slide, l, y + row_h - Emu(int(Pt(1))), w, color=G2)


def simple_table(slide, headers: list[str], rows: list[list[str]], l, t, w, h) -> None:
    tbl_shp = slide.shapes.add_table(len(rows) + 1, len(headers), l, t, w, h)
    tbl     = tbl_shp.table
    col_w   = w // len(headers)
    for ci in range(len(headers)):
        tbl.columns[ci].width = col_w

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        p = cell.text_frame.paragraphs[0]
        if p.runs:
            r = p.runs[0]
            r.font.color.rgb = WHITE
            r.font.bold      = True
            r.font.size      = Pt(9)
            r.font.name      = FONT

    for ri, row in enumerate(rows):
        bg = G1 if ri % 2 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                r = p.runs[0]
                r.font.size      = Pt(10)
                r.font.color.rgb = G5
                r.font.name      = FONT


def col_setup(half_cw: float | None = None):
    """Returns (half, R_L) — left-half width and right-col left edge."""
    half = (CW / 2) - Inches(0.15) if half_cw is None else half_cw
    r_l  = ML + half + Inches(0.35)
    return half, r_l


def draw_divider(slide, body_t, half) -> None:
    vline(slide, ML + half + Inches(0.15), body_t,
          SH - body_t - Inches(0.50), color=G2, name="!!col-divider")


# ══════════════════════════════════════════════════════════════
#  SLIDE 1  —  TITLE
# ══════════════════════════════════════════════════════════════

def build_slide_1(sl) -> None:
    add_morph_furniture(sl, "01")
    eyebrow(sl, "Presentation · Case Studies")

    # Large title
    tb(sl, "Connected", ML, MT + Inches(0.25), Inches(8.5), Inches(0.85),
       size=60, bold=True, name="!!title-line-1")
    tb(sl, "Vehicles",  ML, MT + Inches(1.10), Inches(8.5), Inches(0.85),
       size=60, bold=True, name="!!title-line-2")

    tb(sl,
       "Architecture, real-world deployments, safety data,\n"
       "and what four major case studies actually reveal.",
       ML, MT + Inches(2.10), Inches(7.5), Inches(0.70),
       size=15, color=G3)

    # Tag pills
    tags = ["Tesla Autopilot", "Waymo One", "GM OnStar", "iRASTE Nagpur"]
    tag_w, tag_h = Inches(1.68), Inches(0.28)
    tag_t = MT + Inches(3.10)
    for i, tag in enumerate(tags):
        lx = ML + i * (tag_w + Inches(0.10))
        rect(sl, lx, tag_t, tag_w, tag_h, fill=BLACK)
        tb(sl, tag, lx + Inches(0.10), tag_t + Inches(0.04),
           tag_w - Inches(0.14), tag_h - Inches(0.06),
           size=9, bold=True, color=WHITE)

    hline(sl, ML, SH - Inches(1.10), CW, color=G2)
    tb(sl,
       "Sources: NHTSA EA 22-002  ·  Waymo Safety Hub  ·  INAI/DST iRASTE 2024"
       "  ·  GM Authority  ·  MoRTH AIS-140",
       ML, SH - Inches(0.95), CW, Inches(0.25), size=8, color=G3)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2  —  OVERVIEW
# ══════════════════════════════════════════════════════════════

def build_slide_2(sl) -> None:
    add_morph_furniture(sl, "02")
    eyebrow(sl, "Section 01 — Foundation")
    heading(sl, "What are connected vehicles?")

    BT   = MT + Inches(1.10)
    half, R_L = col_setup()
    draw_divider(sl, BT, half)

    # ── Left ──
    tb(sl,
       "A connected vehicle uses an onboard cellular modem, GPS, and short-range radio "
       "to exchange data with other vehicles, road infrastructure, back-end cloud servers, "
       "and pedestrian devices — in real time, while moving.",
       ML, BT, half, Inches(0.80), size=11, color=G5)

    tb(sl,
       "SAE International defines six autonomy levels. Most cars on public roads today "
       "operate at Level 1–2 (driver assistance). Fully autonomous Level 4 runs only "
       "inside tight geofenced zones in a handful of cities.",
       ML, BT + Inches(0.88), half, Inches(0.70), size=11, color=G5)

    note_box(sl,
             "79% of all new cars sold globally in 2024 shipped with an OEM-embedded "
             "telematics system — up from 75% in 2023.  (Counterpoint Research)",
             ML, BT + Inches(1.66), half, Inches(0.50))

    simple_table(sl,
        ["SAE Level", "Description", "Who drives"],
        [
            ["L0 – L2", "No automation → Partial automation (lane-keep, ACC)", "Human at all times"],
            ["L3",      "Conditional — system drives, human on standby",        "Human when requested"],
            ["L4",      "High automation — no human needed within ODD",         "System (geofenced)"],
            ["L5",      "Full automation — any road, any condition",            "System always"],
        ],
        ML, BT + Inches(2.26), half, Inches(1.55))

    # ── Right ──
    tb(sl, "COMMUNICATION ECOSYSTEM", R_L, BT, half - Inches(0.1),
       Inches(0.22), size=8, bold=True, color=G4)

    cards = [
        ("V2V",      "Vehicle-to-Vehicle",
         "Two cars share speed, position, braking state. Under 10 ms latency with "
         "DSRC (802.11p). First deployed in the 2017 Cadillac CTS.",
         BLACK, WHITE, RGBColor(0xAA, 0xAA, 0xAA)),
        ("V2I",      "Vehicle-to-Infrastructure",
         "Traffic signals, construction warnings, speed advisories. Needs roadside "
         "hardware — the main bottleneck for nationwide rollout.",
         WHITE, BLACK, G4),
        ("V2C / V2X","Vehicle-to-Cloud / Everything",
         "OTA software updates, fleet diagnostics, remote commands. The most "
         "widely deployed communication mode today.",
         G1, BLACK, G4),
    ]
    card_h = Inches(1.28)
    for i, (big, sub, body, bg, tc, sc) in enumerate(cards):
        cy = BT + Inches(0.30) + i * (card_h + Inches(0.09))
        rect(sl, R_L, cy, half - Inches(0.1), card_h, fill=bg, line=G2)
        tb(sl, big,  R_L + Inches(0.14), cy + Inches(0.09),
           Inches(1.6), Inches(0.40), size=22, bold=True, color=tc)
        tb(sl, sub,  R_L + Inches(0.14), cy + Inches(0.47),
           half - Inches(0.35), Inches(0.20), size=8, bold=True, color=sc)
        tb(sl, body, R_L + Inches(0.14), cy + Inches(0.66),
           half - Inches(0.35), Inches(0.55), size=10, color=tc if bg == BLACK else G5)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3  —  TESLA AUTOPILOT
# ══════════════════════════════════════════════════════════════

def build_slide_3(sl) -> None:
    add_morph_furniture(sl, "03")
    eyebrow(sl, "Case Study 01")
    heading(sl, "Tesla Autopilot & Full Self-Driving")

    BT   = MT + Inches(1.10)
    half, R_L = col_setup()
    draw_divider(sl, BT, half)

    # ── Left ──
    tb(sl, "BACKGROUND", ML, BT, half, Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, ML, BT + Inches(0.20), half, color=BLACK)
    tb(sl,
       "Tesla launched Autopilot in October 2014 — the first mass-market ADAS suite. "
       "Full Self-Driving (FSD) is a $15,000 optional package. Both are SAE Level 2: "
       "the driver must supervise at all times.",
       ML, BT + Inches(0.25), half, Inches(0.65), size=10, color=G5)

    tb(sl, "TECHNOLOGY", ML, BT + Inches(0.96), half, Inches(0.20),
       size=8, bold=True, color=G4)
    hline(sl, ML, BT + Inches(1.16), half, color=BLACK)
    tb(sl,
       "Tesla's system uses cameras and neural networks only — no LiDAR. Eight cameras "
       "feed a custom FSD Computer (144 TOPS). All software updates are delivered "
       "over-the-air, making OTA the primary recall mechanism.",
       ML, BT + Inches(1.21), half, Inches(0.65), size=10, color=G5)

    tb(sl, "NHTSA INVESTIGATION — EA 22-002", ML, BT + Inches(1.92),
       half, Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, ML, BT + Inches(2.12), half, color=BLACK)

    sb_w = (half - Inches(0.16)) / 3
    for i, (num, lbl) in enumerate([
        ("956",  "Crashes analyzed\nJan 2018–Aug 2023"),
        ("29",   "Fatalities in\nfull dataset"),
        ("2M+",  "Vehicles recalled\nDec 2023"),
    ]):
        stat_block(sl, num, lbl,
                   ML + i * (sb_w + Inches(0.08)),
                   BT + Inches(2.17), sb_w)

    tb(sl, "Source: NHTSA Engineering Analysis EA 22-002, 2024",
       ML, BT + Inches(3.15), half, Inches(0.20), size=8, color=G3)

    # ── Right ──
    tb(sl, "KEY FINDINGS", R_L, BT, half - Inches(0.1), Inches(0.20),
       size=8, bold=True, color=G4)
    hline(sl, R_L, BT + Inches(0.20), half - Inches(0.1), color=BLACK)

    bullet_list(sl, [
        "At least 14 crashes into first-responder vehicles (fire trucks, police cruisers "
        "with active lights) — 15 injuries, 1 death",
        'NHTSA concluded driver monitoring was "not appropriate for Autopilot\'s permissive '
        "operating capabilities\" — a critical safety gap",
        "Drivers were treating an L2 system as L3 or L4 — hands off, eyes off the road",
    ], R_L, BT + Inches(0.28), half - Inches(0.1), row_h=Inches(0.50), size=10)

    tb(sl, "RECALLS & RESPONSE", R_L, BT + Inches(1.82),
       half - Inches(0.1), Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, R_L, BT + Inches(2.02), half - Inches(0.1), color=BLACK)

    simple_table(sl,
        ["Date", "Scope", "Issue"],
        [
            ["Feb 2023", "363,000 vehicles",    "FSD violated speed limits; intersection errors"],
            ["Dec 2023", "~2,000,000 vehicles", "Autosteer misuse; weak driver attention monitoring"],
        ],
        R_L, BT + Inches(2.07), half - Inches(0.1), Inches(0.84))

    note_box(sl,
             "Lesson: L2 systems need hard, continuous driver monitoring — not just "
             "hands-on-wheel torque sensing. OTA recalls are fast but unverifiable from outside.",
             R_L, BT + Inches(3.06), half - Inches(0.1), Inches(0.55))


# ══════════════════════════════════════════════════════════════
#  SLIDE 4  —  WAYMO ONE
# ══════════════════════════════════════════════════════════════

def build_slide_4(sl) -> None:
    add_morph_furniture(sl, "04")
    eyebrow(sl, "Case Study 02")
    heading(sl, "Waymo One — Commercial Robotaxi")

    BT   = MT + Inches(1.10)
    half, R_L = col_setup()
    draw_divider(sl, BT, half)

    # ── Left ──
    tb(sl, "BACKGROUND", ML, BT, half, Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, ML, BT + Inches(0.20), half, color=BLACK)
    tb(sl,
       "Google's self-driving project started in 2009. It became Waymo (an Alphabet "
       "subsidiary) in 2016. Waymo One launched on December 5, 2018 in Phoenix — "
       "the first paying robotaxi in the US, initially with safety drivers.",
       ML, BT + Inches(0.25), half, Inches(0.65), size=10, color=G5)

    tb(sl, "TECHNOLOGY", ML, BT + Inches(0.96), half, Inches(0.20),
       size=8, bold=True, color=G4)
    hline(sl, ML, BT + Inches(1.16), half, color=BLACK)
    tb(sl,
       "Each vehicle has 29 cameras, 5 LiDAR units, 6 radar sensors, and audio detectors. "
       "HD maps built from prior drives give ground-truth context. All driving decisions "
       "run onboard — no cloud dependency while in motion.",
       ML, BT + Inches(1.21), half, Inches(0.65), size=10, color=G5)

    tb(sl, "WEEKLY RIDES GROWTH", ML, BT + Inches(1.96),
       half, Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, ML, BT + Inches(2.16), half, color=BLACK)

    bars = [
        ("May 2023",   0.022, "10k  rides/week"),
        ("May 2024",   0.110, "50k  rides/week"),
        ("April 2025", 0.555, "250k rides/week"),
        ("Late 2025",  1.000, "450k rides/week"),
    ]
    bar_l   = ML + Inches(1.50)
    bar_w   = half - Inches(2.10)
    bar_h   = Inches(0.22)
    bar_gap = Inches(0.07)
    for i, (lbl, pct, val) in enumerate(bars):
        by = BT + Inches(2.22) + i * (bar_h + bar_gap)
        tb(sl, lbl, ML, by, Inches(1.46), bar_h, size=9, color=G3)
        rect(sl, bar_l, by, bar_w, bar_h, fill=G2)                    # track
        rect(sl, bar_l, by, int(bar_w * pct), bar_h, fill=BLACK)      # fill
        tb(sl, val, bar_l + bar_w + Inches(0.07), by,
           Inches(1.2), bar_h, size=9, bold=True)

    tb(sl, "127 million fully autonomous miles logged through Sep 2025",
       ML, BT + Inches(3.50), half, Inches(0.22), size=8, color=G3)

    # ── Right ──
    tb(sl, "SAFETY OUTCOMES (7.14 M rider-only miles)", R_L, BT,
       half - Inches(0.1), Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, R_L, BT + Inches(0.20), half - Inches(0.1), color=BLACK)
    tb(sl,
       "Peer-reviewed data published in Traffic Injury Prevention (2024), "
       "covering Phoenix and San Francisco:",
       R_L, BT + Inches(0.25), half - Inches(0.1), Inches(0.50), size=10, color=G5)

    sb_w = (half - Inches(0.26)) / 2
    stat_block(sl, "85%", "Fewer injury-reported\ncrashes vs humans",
               R_L, BT + Inches(0.83), sb_w)
    stat_block(sl, "55%", "Fewer police-reported\ncrashes",
               R_L + sb_w + Inches(0.08), BT + Inches(0.83), sb_w)

    note_box(sl,
             "Caveat: Waymo operates within pre-mapped zones in clear weather. Human "
             "benchmarks include all roads and conditions. The safety edge is real — "
             "but operating contexts are not equivalent.",
             R_L, BT + Inches(1.74), half - Inches(0.1), Inches(0.58))

    tb(sl, "REGULATORY TIMELINE", R_L, BT + Inches(2.42),
       half - Inches(0.1), Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, R_L, BT + Inches(2.62), half - Inches(0.1), color=BLACK)

    tl_items = [
        ("DEC 2018", "First commercial launch — Phoenix, with safety drivers"),
        ("NOV 2019", "First rides with no safety driver (limited, Chandler AZ)"),
        ("OCT 2020", "Fully driverless service opened to the general public — 50 sq mi zone"),
        ("AUG 2023", "Full public driverless rides in San Francisco — 4 years after first test"),
        ("2024–25",  "Expanded to Los Angeles, Austin; 3,000 vehicles in fleet"),
    ]
    for i, (yr, ev) in enumerate(tl_items):
        ty = BT + Inches(2.70) + i * Inches(0.37)
        tb(sl, yr, R_L, ty, Inches(0.70), Inches(0.35), size=8, bold=True, color=G3)
        tb(sl, ev, R_L + Inches(0.72), ty, half - Inches(0.82), Inches(0.35),
           size=10, color=G5)
        if i < len(tl_items) - 1:
            hline(sl, R_L, ty + Inches(0.36), half - Inches(0.1), color=G2)

    note_box(sl,
             "Lesson: From \"technology ready\" to \"commercially deployed\" took four years — "
             "not because the engineering was unfinished, but because regulation was.",
             R_L, BT + Inches(4.57), half - Inches(0.1), Inches(0.52))


# ══════════════════════════════════════════════════════════════
#  SLIDE 5  —  GM ONSTAR + iRASTE
# ══════════════════════════════════════════════════════════════

def build_slide_5(sl) -> None:
    add_morph_furniture(sl, "05")
    eyebrow(sl, "Case Studies 03 & 04")
    heading(sl, "GM OnStar & iRASTE Nagpur")

    BT   = MT + Inches(1.10)
    half, R_L = col_setup()
    draw_divider(sl, BT, half)

    # ── Left — GM OnStar ──
    tb(sl, "GM ONSTAR — V2C SINCE 1996", ML, BT, half, Inches(0.20),
       size=8, bold=True, color=G4)
    hline(sl, ML, BT + Inches(0.20), half, color=BLACK)
    tb(sl,
       "OnStar launched on three 1997 Cadillac models with one feature: if an airbag "
       "deployed, the car called a human operator. It is the longest-running connected "
       "vehicle service. Core safety features are free for 8 years on all 2025+ GM "
       "vehicles; premium services (Super Cruise, remote commands) are the paid layer.",
       ML, BT + Inches(0.25), half, Inches(0.80), size=10, color=G5)

    sb_w = (half - Inches(0.08)) / 2
    stat_block(sl, "12M",   "Global subscribers\nend-2025",
               ML, BT + Inches(1.10), sb_w)
    stat_block(sl, "3,500", "Auto crash responses\nper month",
               ML + sb_w + Inches(0.08), BT + Inches(1.10), sb_w)

    bullet_list(sl, [
        "Injury Severity Prediction uses Delta-V, seatbelt status, and impact direction "
        "to flag high-risk crashes for priority dispatch",
        "Super Cruise (hands-free highway driving): 625,000+ active users; ~$200M revenue",
        "2025 Texas floods: 4,000 crisis calls answered, 12+ direct rescues supported",
    ], ML, BT + Inches(2.08), half, row_h=Inches(0.42), size=10)

    note_box(sl,
             "Lesson: Connectivity as a safety layer drives decades-long subscriber loyalty. "
             "$5.4 billion in deferred revenue (+65% YoY) proves it.  (GM Authority, 2026)",
             ML, BT + Inches(3.42), half, Inches(0.52))

    # ── Right — iRASTE ──
    tb(sl, "iRASTE — NAGPUR, INDIA (SEPT 2021)", R_L, BT,
       half - Inches(0.1), Inches(0.20), size=8, bold=True, color=G4)
    hline(sl, R_L, BT + Inches(0.20), half - Inches(0.1), color=BLACK)
    tb(sl,
       "Partners: Nagpur Municipal Corporation, CSIR-CRRI, Mahindra, IIIT-Hyderabad, Intel. "
       "AI-based ADAS on NMC public buses, RFID at 10 high-risk intersections, driver "
       "behaviour monitoring, and predictive AI to identify grey spots before they "
       "become accident hotspots.",
       R_L, BT + Inches(0.25), half - Inches(0.1), Inches(0.85), size=10, color=G5)

    simple_table(sl,
        ["Metric", "Result"],
        [
            ["Crash reduction (monitored buses vs baseline)", "33%"],
            ["Accidents: ADAS-equipped vs non-ADAS buses",   "41% lower"],
            ["Fatalities in ADAS buses (study period)",      "Zero"],
            ["Lives saved Jul 2023 – Apr 2024",              "36"],
            ["Signal adherence at RFID intersections",       "+24%"],
        ],
        R_L, BT + Inches(1.16), half - Inches(0.1), Inches(1.75))

    note_box(sl,
             "India context: AIS-140 (MoRTH, 2019) mandates GPS/NavIC tracking and SOS "
             "buttons on all commercial vehicles. India connected car market: "
             "$51.6M (2022) → projected $205.6M (2027).  (Frost & Sullivan)",
             R_L, BT + Inches(3.03), half - Inches(0.1), Inches(0.65))

    tb(sl, "Source: INAI/DST iRASTE Detailed Report, 2024",
       R_L, BT + Inches(3.78), half - Inches(0.1), Inches(0.22), size=8, color=G3)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6  —  CONCLUSION
# ══════════════════════════════════════════════════════════════

def build_slide_6(sl) -> None:
    add_morph_furniture(sl, "06")
    eyebrow(sl, "Conclusion")
    heading(sl, "Key takeaways")

    BT     = MT + Inches(1.05)
    card_w = (CW - Inches(0.30)) / 3
    card_h = Inches(2.35)

    lessons = [
        ("01", "Safety works — with limits",
         "Waymo's 85% crash reduction and iRASTE's 41% accident drop are backed by "
         "peer-reviewed data. But both operate in controlled, pre-defined environments. "
         "Scaling to all roads, all weather, all conditions — that gap is still open."),
        ("02", "Systems need honest design",
         "Tesla's core problem wasn't the cameras — it was naming a Level 2 system "
         '"Autopilot" and shipping weak driver monitoring. What drivers believe about '
         "a system shapes how they use it. Marketing language is a safety decision."),
        ("03", "Infrastructure is the real gap",
         "Ford committed to V2X hardware in all new US vehicles from 2022. But V2X "
         "only works when roads are also equipped. The hardware is ready; the "
         "infrastructure investment has not followed. This is true globally."),
    ]

    for i, (num, title, body) in enumerate(lessons):
        lx = ML + i * (card_w + Inches(0.15))
        rect(sl, lx, BT, card_w, card_h, fill=G1, line=G2)
        tb(sl, num,   lx + Inches(0.15), BT + Inches(0.10),
           card_w - Inches(0.20), Inches(0.55), size=42, bold=True, color=G2)
        tb(sl, title, lx + Inches(0.15), BT + Inches(0.62),
           card_w - Inches(0.25), Inches(0.36), size=9, bold=True, color=G4)
        tb(sl, body,  lx + Inches(0.15), BT + Inches(1.00),
           card_w - Inches(0.25), Inches(1.25), size=10, color=G5)

    # Bottom section
    DIV_T = BT + card_h + Inches(0.16)
    hline(sl, ML, DIV_T, CW, color=G2)

    half, R_L = col_setup()
    FT = DIV_T + Inches(0.14)

    tb(sl, "FUTURE DIRECTIONS", ML, FT, half, Inches(0.20),
       size=8, bold=True, color=G4)
    bullet_list(sl, [
        "5G-V2X enabling sub-millisecond latency for collision avoidance at scale",
        "Federated learning — vehicles train shared models without sharing raw data",
        "Vehicle-to-Grid (V2G): EVs as distributed energy storage (GM + PG&E pilot, 2025)",
    ], ML, FT + Inches(0.24), half, row_h=Inches(0.36), size=10)

    vline(sl, ML + half + Inches(0.15), FT,
          SH - FT - Inches(0.48), color=G2)

    tb(sl, "PERSISTENT CHALLENGES", R_L, FT, half - Inches(0.1), Inches(0.20),
       size=8, bold=True, color=G4)
    bullet_list(sl, [
        "Edge cases — unusual pedestrian behaviour, unmapped roads, adverse weather",
        "Regulation lags deployment by 3–5 years in most jurisdictions",
        "Cybersecurity: a connected vehicle is an attack surface with wheels",
    ], R_L, FT + Inches(0.24), half - Inches(0.1), row_h=Inches(0.36), size=10)

    tb(sl,
       "Sources: NHTSA EA 22-002  ·  Waymo Safety Hub / Traffic Injury Prevention 2024  ·  "
       "INAI/DST iRASTE 2024  ·  GM Authority Jan 2026  ·  MoRTH AIS-140  ·  "
       "Counterpoint Research 2024",
       ML, SH - Inches(0.45), CW, Inches(0.24), size=7, color=G3)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════

def main() -> None:
    prs = new_presentation()

    builders = [
        ("01", build_slide_1),
        ("02", build_slide_2),
        ("03", build_slide_3),
        ("04", build_slide_4),
        ("05", build_slide_5),
        ("06", build_slide_6),
    ]

    for i, (num, build) in enumerate(builders):
        sl = blank_slide(prs)
        build(sl)
        if i > 0:                       # no transition on the first slide
            add_morph_transition(sl)

    out_path = "connected_vehicles.pptx"
    prs.save(out_path)

    print(f"\n  Saved → {out_path}")
    print("  Slide size : 13.333\" × 7.5\"  (16:9 widescreen)")
    print("  Transitions: Morph (requires PowerPoint 2019 / Microsoft 365)")
    print("\n  Tip: Shapes named !!watermark, !!accent-line, !!dot")
    print("       animate between slides via Morph — verify in PowerPoint.")


if __name__ == "__main__":
    main()
